"""Audit exercise slides: fenced blocks vs body text vs code boxes."""
import importlib.util
import re
import sys
from pathlib import Path

sys.path.insert(0, "scripts")
from marp_tables import extract_fenced_code_blocks, split_marp_slides
from pptx import Presentation

spec = importlib.util.spec_from_file_location(
    "c", Path("scripts/clean-editable-pptx.py")
)
c = importlib.util.module_from_spec(spec)
assert spec.loader is not None
spec.loader.exec_module(c)

md = split_marp_slides(
    Path("slides/course-complete-marp.md").read_text(encoding="utf-8")
)
prs = Presentation("slides/course-complete-marp-editable-clean.pptx")

search_from = 0
missing_in_pptx = []
body_instead_of_code = []
wrong_placement = []
extra_code = []

for md_index, slide_md in enumerate(md, start=1):
    if not re.search(r"^## Exercise", slide_md, re.MULTILINE):
        continue
    blocks = extract_fenced_code_blocks(slide_md)
    if not blocks:
        continue

    slide_index = c._find_pptx_slide_index(prs, slide_md, start_index=search_from)
    if slide_index is None:
        missing_in_pptx.append((md_index, blocks[0][:50]))
        continue
    search_from = slide_index + 1
    slide = prs.slides[slide_index]
    heading = c._pptx_slide_heading(slide) or ""

    code_shapes = [
        s
        for s in slide.shapes
        if c._is_code_shape(s) and not c._is_inline_monospace_shape(s)
    ]
    if len(code_shapes) != len(blocks):
        extra_code.append(
            (slide_index + 1, len(blocks), len(code_shapes), heading)
        )

    step_numbers = c._code_block_step_numbers(slide_md)

    for bi, block in enumerate(blocks):
        block_norm = c._normalize_text(block)
        if len(block_norm) < 15:
            continue

        # Is block content in a code shape?
        in_code = any(
            c._code_block_match_score(c._shape_text(s), block) >= 8 for s in code_shapes
        )
        if not in_code:
            body_instead_of_code.append(
                (slide_index + 1, heading, bi, block.splitlines()[0][:55])
            )
            continue

        # Step placement check
        if bi < len(step_numbers) and step_numbers[bi] is not None:
            step_n = step_numbers[bi]
            step_shape = c._find_step_label_shape(slide, step_n)
            if step_shape is None:
                wrong_placement.append(
                    (slide_index + 1, heading, f"no Step {step_n} label")
                )
                continue
            # find matching code shape
            match_shape = None
            best = 0
            for s in code_shapes:
                score = c._code_block_match_score(c._shape_text(s), block)
                if score > best:
                    best = score
                    match_shape = s
            if match_shape is None:
                continue
            expected_top = int(step_shape.top) + int(step_shape.height)
            actual_top = int(match_shape.top)
            if actual_top < expected_top - 50000 or actual_top > expected_top + 800000:
                wrong_placement.append(
                    (
                        slide_index + 1,
                        heading,
                        f"block {bi+1} step {step_n} top={actual_top} expected~{expected_top}",
                    )
                )

        # Block content duplicated in body?
        for shape in slide.shapes:
            if not shape.has_text_frame or c._is_chrome_shape(shape):
                continue
            if c._is_code_shape(shape):
                continue
            body = c._normalize_text(c._shape_text(shape))
            if len(body) < 20:
                continue
            if block_norm in body or (
                len(body) > 30 and body in block_norm
            ):
                body_instead_of_code.append(
                    (
                        slide_index + 1,
                        heading,
                        bi,
                        f"DUPE in body: {body[:50]}",
                    )
                )

print("exercise slides with block/box count mismatch:", len(extra_code))
for row in extra_code:
    print(" ", row)

print("\nexercise blocks not in code shapes:", len(body_instead_of_code))
for row in body_instead_of_code[:30]:
    print(" ", row)
if len(body_instead_of_code) > 30:
    print(f"  ... +{len(body_instead_of_code)-30} more")

print("\nstep placement issues:", len(wrong_placement))
for row in wrong_placement[:25]:
    print(" ", row)
