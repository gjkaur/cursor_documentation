import importlib.util
import re
import sys
from pathlib import Path

sys.path.insert(0, "scripts")
from marp_tables import extract_fenced_code_blocks, split_marp_slides
from pptx import Presentation

spec = importlib.util.spec_from_file_location(
    "c", Path("scripts/clean-editable-pptx.py")
)
c = importlib.util.module_from_spec(spec)
assert spec.loader is not None
spec.loader.exec_module(c)

md = split_marp_slides(
    Path("slides/course-complete-marp.md").read_text(encoding="utf-8")
)
prs = Presentation("slides/course-complete-marp-editable-clean.pptx")


def show(idx: int) -> None:
    slide_md = md[idx - 1]
    blocks = extract_fenced_code_blocks(slide_md)
    print(
        "=== pptx",
        idx,
        "md blocks",
        len(blocks),
        "heading",
        (c._markdown_slide_heading(slide_md) or "")[:45],
    )
    slide = prs.slides[idx - 1]
    shapes = sorted(
        [
            shape
            for shape in slide.shapes
            if shape.has_text_frame
            and c._shape_text(shape)
            and not c._is_chrome_shape(shape)
        ],
        key=lambda shape: (int(shape.top), int(shape.left)),
    )
    for shape in shapes:
        text = c._shape_text(shape)
        if c._is_heading_shape(shape):
            kind = "H"
        elif c._is_code_shape(shape):
            kind = "C" + ("F" if c._shape_has_code_fill(shape) else "-")
        else:
            kind = "B"
        preview = text[:72].replace("\n", " | ")
        print(f" {kind} {preview!r}")
    for index, block in enumerate(blocks, start=1):
        preview = block[:72].replace("\n", " | ")
        print(f" MD{index}: {preview!r}")


for slide_index in [60, 62, 64, 66, 72, 82, 86, 88, 125, 127]:
    show(slide_index)

# count exercise slides where body contains prompt-like fenced first line
search_from = 0
problems = []
for md_index, slide_md in enumerate(md, start=1):
    if not re.search(r"^## Exercise", slide_md, re.MULTILINE):
        continue
    blocks = extract_fenced_code_blocks(slide_md)
    if not blocks:
        continue
    slide_index = c._find_pptx_slide_index(prs, slide_md, start_index=search_from)
    if slide_index is None:
        continue
    search_from = slide_index + 1
    slide = prs.slides[slide_index]
    code_text = "\n".join(
        c._shape_text(shape)
        for shape in slide.shapes
        if c._is_code_shape(shape) and not c._is_inline_monospace_shape(shape)
    )
    body_text = "\n".join(
        c._shape_text(shape)
        for shape in slide.shapes
        if shape.has_text_frame
        and not c._is_chrome_shape(shape)
        and not c._is_heading_shape(shape)
        and not c._is_code_shape(shape)
        and c._shape_text(shape)
    )
    for block in blocks:
        first = block.splitlines()[0].strip()
        norm_first = c._normalize_text(first)
        norm_code = c._normalize_text(code_text)
        norm_body = c._normalize_text(body_text)
        in_code = norm_first in norm_code or c._normalize_text(block) in norm_code
        in_body = norm_first in norm_body or c._normalize_text(block)[:80] in norm_body
        if in_body and not in_code:
            problems.append((slide_index + 1, first[:50]))
        elif not in_code and not in_body:
            problems.append((slide_index + 1, "MISSING:" + first[:40]))

print("\nExercise fenced content in body or missing:", len(problems))
for row in problems[:30]:
    print(" ", row)
