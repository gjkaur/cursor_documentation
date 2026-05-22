import importlib.util
import re
import sys
from pathlib import Path

sys.path.insert(0, "scripts")
from marp_tables import extract_fenced_code_blocks, split_marp_slides
from pptx import Presentation

spec = importlib.util.spec_from_file_location(
    "c", Path("scripts/clean-editable-pptx.py")
)
c = importlib.util.module_from_spec(spec)
assert spec.loader is not None
spec.loader.exec_module(c)

md = split_marp_slides(
    Path("slides/course-complete-marp.md").read_text(encoding="utf-8")
)
prs = Presentation("slides/course-complete-marp-editable-clean.pptx")
search_from = 0
body_leaks = []

for md_index, slide_md in enumerate(md, start=1):
    blocks = extract_fenced_code_blocks(slide_md)
    if not blocks:
        continue
    slide_index = c._find_pptx_slide_index(prs, slide_md, start_index=search_from)
    if slide_index is None:
        continue
    search_from = slide_index + 1
    slide = prs.slides[slide_index]
    code_norm = c._normalize_text(
        "\n".join(
            c._shape_text(shape)
            for shape in slide.shapes
            if c._is_code_shape(shape) and not c._is_inline_monospace_shape(shape)
        )
    )
    for shape in slide.shapes:
        if not shape.has_text_frame or c._is_chrome_shape(shape) or c._is_heading_shape(shape):
            continue
        if c._is_code_shape(shape):
            continue
        body = c._normalize_text(c._shape_text(shape))
        if not body:
            continue
        for block in blocks:
            block_norm = c._normalize_text(block)
            if len(block_norm) < 20:
                continue
            if block_norm in body or body in block_norm:
                overlap = min(len(block_norm), len(body))
                if overlap > 30:
                    body_leaks.append(
                        (
                            slide_index + 1,
                            c._pptx_slide_heading(slide) or "",
                            block.splitlines()[0][:45],
                            body[:60],
                        )
                    )

print("body shapes overlapping fenced blocks:", len(body_leaks))
for row in body_leaks[:20]:
    print(row)

# exercise slides: count code boxes vs blocks
search_from = 0
short = []
for md_index, slide_md in enumerate(md, start=1):
    if not re.search(r"^## Exercise", slide_md, re.MULTILINE):
        continue
    blocks = extract_fenced_code_blocks(slide_md)
    if not blocks:
        continue
    slide_index = c._find_pptx_slide_index(prs, slide_md, start_index=search_from)
    if slide_index is None:
        continue
    search_from = slide_index + 1
    slide = prs.slides[slide_index]
    codes = [
        s
        for s in slide.shapes
        if c._is_code_shape(s) and not c._is_inline_monospace_shape(s)
    ]
    if len(codes) != len(blocks):
        short.append((slide_index + 1, len(blocks), len(codes), c._pptx_slide_heading(slide)))

print("\nexercise slides block count mismatch:", len(short))
for row in short:
    print(row)
