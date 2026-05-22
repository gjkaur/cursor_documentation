import importlib.util
import re
import sys
from pathlib import Path

sys.path.insert(0, "scripts")
from marp_tables import extract_fenced_code_blocks, split_marp_slides
from pptx import Presentation

spec = importlib.util.spec_from_file_location(
    "c", Path("scripts/clean-editable-pptx.py")
)
c = importlib.util.module_from_spec(spec)
assert spec.loader is not None
spec.loader.exec_module(c)

md = split_marp_slides(
    Path("slides/course-complete-marp.md").read_text(encoding="utf-8")
)
prs = Presentation("slides/course-complete-marp-editable-clean.pptx")
targets = {190, 271, 276, 307, 314, 315, 347}
search_from = 0

for slide_md in md:
    slide_index = c._find_pptx_slide_index(prs, slide_md, start_index=search_from)
    if slide_index is None:
        continue
    search_from = slide_index + 1
    if slide_index + 1 not in targets:
        continue

    blocks = extract_fenced_code_blocks(slide_md)
    slide = prs.slides[slide_index]
    print("=== slide", slide_index + 1, c._pptx_slide_heading(slide), "===")
    print("MD blocks:", len(blocks))
    for i, block in enumerate(blocks):
        print("  block", i, repr(block[:100]))
    print("Shapes:")
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = c._shape_text(shape)[:120].replace("\n", " | ")
        is_code = c._is_code_shape(shape)
        inline = c._is_inline_monospace_shape(shape)
        heading = c._is_heading_shape(shape)
        matches = c._code_shape_matches_any_block(c._shape_text(shape), blocks)
        print(
            f"  code={is_code} inline={inline} heading={heading} "
            f"matches_block={matches} text={text!r}"
        )
    print()
