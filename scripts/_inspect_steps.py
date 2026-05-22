import importlib.util
import re
import sys
from pathlib import Path

sys.path.insert(0, "scripts")
from marp_tables import extract_fenced_code_blocks, split_marp_slides
from pptx import Presentation

spec = importlib.util.spec_from_file_location(
    "c", Path("scripts/clean-editable-pptx.py")
)
c = importlib.util.module_from_spec(spec)
assert spec.loader is not None
spec.loader.exec_module(c)

md = split_marp_slides(
    Path("slides/course-complete-marp.md").read_text(encoding="utf-8")
)
prs = Presentation("slides/course-complete-marp-editable-clean.pptx")

targets = [62, 67, 86, 190, 243, 244]
search_from = 0
for slide_md in md:
    slide_index = c._find_pptx_slide_index(prs, slide_md, start_index=search_from)
    if slide_index is None:
        continue
    search_from = slide_index + 1
    if slide_index + 1 not in targets:
        continue
    slide = prs.slides[slide_index]
    blocks = extract_fenced_code_blocks(slide_md)
    steps = c._code_block_step_numbers(slide_md)
    print("===", slide_index + 1, c._pptx_slide_heading(slide), "===")
    print("blocks/steps:", list(zip(steps, [b.splitlines()[0][:40] for b in blocks])))
    for shape in sorted(
        [s for s in slide.shapes if s.has_text_frame],
        key=lambda s: (int(s.top), int(s.left)),
    ):
        if c._is_chrome_shape(shape):
            continue
        kind = "CODE" if c._is_code_shape(shape) else ("HDR" if c._is_heading_shape(shape) else "body")
        print(f"  y={int(shape.top)/914400:.2f}in {kind}: {c._shape_text(shape)[:90]!r}")
    print()
