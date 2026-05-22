#!/usr/bin/env python3
"""Audit fenced code blocks in markdown vs styled code boxes in PPTX."""

from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from marp_tables import extract_fenced_code_blocks, split_marp_slides
from pptx import Presentation

_spec = importlib.util.spec_from_file_location(
    "clean_editable_pptx",
    Path(__file__).resolve().parent / "clean-editable-pptx.py",
)
c = importlib.util.module_from_spec(_spec)
assert _spec.loader is not None
_spec.loader.exec_module(c)


def audit(md_path: Path, pptx_path: Path) -> None:
    markdown = md_path.read_text(encoding="utf-8")
    md_slides = split_marp_slides(markdown)
    prs = Presentation(str(pptx_path))
    search_from = 0
    missing: list[tuple[int, int, str, str]] = []
    expected = 0
    found = 0

    for md_index, slide_md in enumerate(md_slides, start=1):
        blocks = extract_fenced_code_blocks(slide_md)
        if not blocks:
            continue
        slide_index = c._find_pptx_slide_index(prs, slide_md, start_index=search_from)
        if slide_index is None:
            for block in blocks:
                missing.append((md_index, 0, block.splitlines()[0][:50], "slide not matched"))
            expected += len(blocks)
            continue
        search_from = slide_index + 1
        slide = prs.slides[slide_index]
        code_shapes = [
            shape
            for shape in slide.shapes
            if c._is_code_shape(shape) and not c._is_inline_monospace_shape(shape)
        ]
        expected += len(blocks)
        found += min(len(blocks), len(code_shapes))

        if len(code_shapes) < len(blocks):
            heading = c._pptx_slide_heading(slide) or ""
            for block in blocks[len(code_shapes) :]:
                missing.append(
                    (
                        md_index,
                        slide_index + 1,
                        block.splitlines()[0][:60],
                        heading[:40],
                    )
                )

    print(f"{pptx_path.name}: expected {expected} fenced blocks, styled boxes ~{found}, missing {len(missing)}")
    for item in missing[:25]:
        print(f"  md slide {item[0]} -> pptx {item[1]} | {item[3]!r} | {item[2]!r}")
    if len(missing) > 25:
        print(f"  ... and {len(missing) - 25} more")


if __name__ == "__main__":
    audit(
        Path("slides/course-complete-marp.md"),
        Path("slides/course-complete-marp-editable-clean.pptx"),
    )
