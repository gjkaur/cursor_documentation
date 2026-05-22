#!/usr/bin/env python3
"""Quick audit of exported PPTX formatting consistency."""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

import importlib.util

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

_spec = importlib.util.spec_from_file_location(
    "clean_editable_pptx",
    Path(__file__).resolve().parent / "clean-editable-pptx.py",
)
c = importlib.util.module_from_spec(_spec)
assert _spec.loader is not None
_spec.loader.exec_module(c)


def audit(path: Path) -> dict:
    prs = Presentation(str(path))
    slide_height = int(prs.slide_height)
    body_wrong: list[tuple[int, float, str]] = []
    heading_wrong: list[tuple[int, float, str]] = []
    overlaps: list[int] = []
    footer_viol: list[int] = []
    summary_no_table: list[int] = []
    code_wrong: list[tuple[int, float, str]] = []

    for index, slide in enumerate(prs.slides, 1):
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame and c._is_heading_shape(shape):
                title = c._shape_text(shape)[:60]
                break

        if c._slide_has_layout_issues(slide, slide_height):
            footer_line = int(slide_height * c.FOOTER_SAFE_RATIO)
            shapes = c._content_layout_shapes(slide)
            if any(c._shape_bounds(shape)[3] > footer_line for shape in shapes):
                footer_viol.append(index)
            else:
                overlaps.append(index)

        if title == "Module Summary":
            has_table = any(
                shape.shape_type == MSO_SHAPE_TYPE.TABLE for shape in slide.shapes
            )
            if not has_table:
                summary_no_table.append(index)

        for shape in slide.shapes:
            if c._is_chrome_shape(shape):
                continue
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                continue
            if not shape.has_text_frame:
                continue
            text = c._shape_text(shape)
            if not text:
                continue
            max_pt = c._max_run_font_pt(shape)
            if max_pt is None:
                continue
            if c._is_code_shape(shape):
                if abs(max_pt - c.CODE_FONT_PT) > 0.5:
                    code_wrong.append((index, round(max_pt, 1), text[:40]))
            elif c._is_inline_monospace_shape(shape):
                if abs(max_pt - c.MIN_BODY_FONT_PT) > 0.5:
                    body_wrong.append((index, round(max_pt, 1), text[:40]))
            elif c._is_heading_shape(shape):
                if c._is_lead_slide(slide) or c._is_lesson_lead_slide(slide):
                    continue
                if abs(max_pt - c.HEADING_FONT_PT) > 0.5:
                    heading_wrong.append((index, round(max_pt, 1), text[:40]))
            else:
                if abs(max_pt - c.MIN_BODY_FONT_PT) > 0.5:
                    body_wrong.append((index, round(max_pt, 1), text[:40]))

    return {
        "slides": len(prs.slides),
        "footer_viol": footer_viol,
        "overlaps": overlaps,
        "body_wrong": body_wrong,
        "heading_wrong": heading_wrong,
        "code_wrong": code_wrong,
        "summary_no_table": summary_no_table,
    }


def print_report(path: Path, result: dict) -> None:
    print(f"=== {path.name} ({result['slides']} slides) ===")
    print(f"  footer violations: {len(result['footer_viol'])} {result['footer_viol'][:10]}")
    print(f"  overlaps: {len(result['overlaps'])} {result['overlaps'][:10]}")
    print(f"  body font != {int(c.MIN_BODY_FONT_PT)}pt: {len(result['body_wrong'])}")
    for item in result["body_wrong"][:5]:
        print(f"    slide {item[0]}: {item[1]}pt \"{item[2]}\"")
    print(f"  heading font != 32pt: {len(result['heading_wrong'])}")
    for item in result["heading_wrong"][:5]:
        print(f"    slide {item[0]}: {item[1]}pt \"{item[2]}\"")
    print(f"  code font issues: {len(result['code_wrong'])}")
    for item in result["code_wrong"][:5]:
        print(f"    slide {item[0]}: {item[1]}pt \"{item[2]}\"")
    print(f"  module summary without table: {result['summary_no_table']}")
    print()


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "paths",
        nargs="*",
        type=Path,
        help="PPTX files to audit (default: course + all modules)",
    )
    args = parser.parse_args()

    if args.paths:
        paths = args.paths
    else:
        paths = [Path("slides/course-complete-marp-editable-clean.pptx")]
        paths += sorted(Path("slides").glob("module-*-marp-editable-clean.pptx"))

    total_issues = 0
    for path in paths:
        if not path.exists():
            print(f"Missing: {path}")
            continue
        result = audit(path)
        print_report(path, result)
        total_issues += (
            len(result["footer_viol"])
            + len(result["overlaps"])
            + len(result["body_wrong"])
            + len(result["heading_wrong"])
            + len(result["code_wrong"])
            + len(result["summary_no_table"])
        )

    return 1 if total_issues else 0


if __name__ == "__main__":
    raise SystemExit(main())
