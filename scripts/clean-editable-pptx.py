#!/usr/bin/env python3
"""
Clean Marp --pptx-editable exports for manual editing in PowerPoint.

LibreOffice often decomposes Gaia's diagonal background into hundreds of tiny
shapes. This script keeps text boxes, sets a flat Gaia-style background, and
removes decorative/non-text shapes.
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt

from marp_tables import extract_fenced_code_blocks, extract_markdown_tables, split_marp_slides


GAIA_CREAM = RGBColor(0xFF, 0xFF, 0xFF)
HEADING_RED = RGBColor(0xCC, 0x00, 0x00)
BODY_BLACK = RGBColor(0x00, 0x00, 0x00)
TABLE_HEADER_BG = RGBColor(0xFF, 0xEB, 0xEB)
TABLE_ROW_BG = RGBColor(0xF7, 0xF7, 0xF7)
TABLE_ROW_ALT_BG = RGBColor(0xFF, 0xFF, 0xFF)
TABLE_BORDER = "E0E0E0"
BODY_FONT = "Calibri"
CODE_FONT = "Consolas"
CODE_FONT_PT = 18
CODE_BG = RGBColor(0xF4, 0xF4, 0xF4)
CODE_BORDER = RGBColor(0xD0, 0xD0, 0xD0)
CODE_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
CODE_INSET = Emu(20000)
CODE_LINE_HEIGHT_PT = 20
CODE_CHAR_WIDTH_PT = 10.8
CODE_SIZE_BUFFER_PT = 2
MIN_BODY_FONT_PT = 20
BODY_LINE_SPACING = 1.5
HEADING_FONT_PT = 32
LEAD_H1_FONT_PT = 40
LEAD_H2_FONT_PT = 32
TABLE_FONT_PT = 18
TABLE_LINE_SPACING = 1.0
TABLE_ROW_HEIGHT_PT = int(round(TABLE_FONT_PT * TABLE_LINE_SPACING))
TABLE_CHAR_WIDTH_PT = 8.4
TABLE_CELL_INSET = Emu(30000)
TABLE_SIZE_BUFFER_PT = 6
TABLE_MIN_COL_WIDTH = Emu(520000)
TABLE_MAX_COL_CHARS = 44
CHROME_MAX_FONT_PT = 15
HEADING_MIN_FONT_PT = 26
MIN_LAYOUT_FONT_PT = 14
OVERLAP_GAP = Emu(25000)
CONTENT_GAP = Emu(80000)
HEADING_BODY_GAP = Emu(300000)
FOOTER_SAFE_RATIO = 0.86
CONTENT_LEFT_RATIO = 0.061
HEADING_TOP_CONTENT = Emu(850000)
SECTION_CENTER_TOLERANCE = 0.08
MARP_HEADER_RE = re.compile(r"<!--\s*_header:\s*['\"](.+?)['\"]\s*-->")
MARP_FOOTER_RE = re.compile(r"<!--\s*_footer:\s*['\"](.+?)['\"]\s*-->")
CHROME_FONT_PT = 13.5
HEADER_CHROME_TOP_RATIO = 0.15
FOOTER_CHROME_TOP_RATIO = 0.85


def _shape_fill_rgb(shape) -> str | None:
    try:
        fill = shape.fill
        if fill.type is None or fill.type != 1:
            return None
        return str(fill.fore_color.rgb).upper()
    except Exception:
        return None


def _is_full_slide(shape, slide_width: int, slide_height: int, tolerance: float = 0.98) -> bool:
    w, h = int(shape.width), int(shape.height)
    return w >= slide_width * tolerance and h >= slide_height * tolerance


def _shape_text(shape) -> str:
    try:
        if not shape.has_text_frame:
            return ""
        return shape.text_frame.text.strip()
    except Exception:
        return ""


def _set_slide_background(slide, color: RGBColor = GAIA_CREAM) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _max_run_font_pt(shape) -> float | None:
    if not shape.has_text_frame:
        return None
    sizes = [
        run.font.size.pt
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.font.size is not None
    ]
    return max(sizes) if sizes else None


def _is_chrome_shape(shape) -> bool:
    """Slide chrome only — module header, footer tagline, page number.

    Do not use font size alone: LibreOffice exports table cells around 14pt and
    those must remain matchable for native table rebuild.
    """
    text = _shape_text(shape)
    if not text:
        return False
    stripped = text.strip()
    if stripped.isdigit() and len(stripped) <= 4:
        return True
    if stripped.startswith("Module ") and " — " in stripped:
        return True
    if "Cursor Training Program" in stripped:
        return True
    if stripped.startswith("Springpeople"):
        return True
    return False


def _scale_body_fonts(shape, min_pt: float = MIN_BODY_FONT_PT) -> None:
    _normalize_body_fonts(shape, min_pt=min_pt)


def _normalize_body_fonts(shape, min_pt: float = MIN_BODY_FONT_PT) -> None:
    if not shape.has_text_frame or _is_chrome_shape(shape):
        return
    if _is_heading_shape(shape):
        return
    if _is_code_shape(shape):
        return
    if _shape_uses_monospace(shape):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = CODE_FONT
                run.font.size = Pt(min_pt)
                run.font.color.rgb = CODE_TEXT
        return
    for paragraph in shape.text_frame.paragraphs:
        _apply_body_paragraph_spacing(paragraph)
        for run in paragraph.runs:
            run.font.size = Pt(min_pt)
            run.font.color.rgb = BODY_BLACK
    _apply_body_list_formatting(shape)


def _normalize_table_fonts(shape) -> None:
    if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
        return
    table = shape.table
    for row_idx in range(len(table.rows)):
        for col_idx in range(len(table.columns)):
            cell = table.cell(row_idx, col_idx)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.space_before = Pt(0)
                paragraph.space_after = Pt(0)
                paragraph.line_spacing = TABLE_LINE_SPACING
                for run in paragraph.runs:
                    run.font.size = Pt(TABLE_FONT_PT)
                    if row_idx == 0:
                        run.font.bold = True
                        run.font.color.rgb = HEADING_RED
                    else:
                        run.font.color.rgb = BODY_BLACK


def _normalize_code_block_fonts(shape) -> None:
    if not _is_code_shape(shape) or not shape.has_text_frame:
        return
    if not _shape_has_code_fill(shape):
        shape.fill.solid()
        shape.fill.fore_color.rgb = CODE_BG
        shape.line.color.rgb = CODE_BORDER
        shape.line.width = Pt(0.75)
    frame = shape.text_frame
    frame.margin_left = CODE_INSET
    frame.margin_right = CODE_INSET
    frame.margin_top = CODE_INSET
    frame.margin_bottom = CODE_INSET
    for paragraph in frame.paragraphs:
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = Pt(CODE_LINE_HEIGHT_PT)
        for run in paragraph.runs:
            run.font.name = CODE_FONT
            run.font.size = Pt(CODE_FONT_PT)
            run.font.color.rgb = CODE_TEXT
    line_count = max(1, len(_code_lines(_shape_text(shape))))
    frame_pad = int(CODE_INSET) * 2
    buffer = int(Pt(CODE_SIZE_BUFFER_PT).emu)
    line_height = int(Pt(CODE_LINE_HEIGHT_PT).emu)
    shape.height = line_count * line_height + frame_pad + buffer


def _shape_text_color(shape) -> RGBColor:
    text = _shape_text(shape)
    if text.startswith("Module ") and " — " in text:
        return HEADING_RED
    max_pt = _max_run_font_pt(shape)
    if max_pt is not None and max_pt >= HEADING_MIN_FONT_PT:
        return HEADING_RED
    return BODY_BLACK


def _apply_text_colors(shape) -> None:
    if not shape.has_text_frame or _shape_uses_monospace(shape):
        return
    color = _shape_text_color(shape)
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = color


def _is_heading_shape(shape) -> bool:
    if not shape.has_text_frame or _is_chrome_shape(shape) or _is_code_shape(shape):
        return False
    text = _shape_text(shape)
    if text.startswith("Cursor Training Program"):
        return False
    max_pt = _max_run_font_pt(shape)
    return max_pt is not None and max_pt >= HEADING_MIN_FONT_PT


def _parse_marp_frontmatter_chrome(markdown: str) -> tuple[str | None, str | None]:
    text = markdown.lstrip("\ufeff")
    if not text.startswith("---"):
        return None, None
    end = text.find("\n---\n", 3)
    if end == -1:
        return None, None
    header = footer = None
    for line in text[3:end].splitlines():
        stripped = line.strip()
        if stripped.startswith("header:"):
            header = stripped.split(":", 1)[1].strip().strip("'\"")
        elif stripped.startswith("footer:"):
            footer = stripped.split(":", 1)[1].strip().strip("'\"")
    return header, footer


def _cumulative_slide_chrome(
    md_slides: list[str],
    default_header: str | None,
    default_footer: str | None,
) -> list[tuple[str, str]]:
    header = default_header or ""
    footer = default_footer or ""
    chrome: list[tuple[str, str]] = []
    for slide_md in md_slides:
        header_match = MARP_HEADER_RE.search(slide_md)
        if header_match:
            header = header_match.group(1)
        footer_match = MARP_FOOTER_RE.search(slide_md)
        if footer_match:
            footer = footer_match.group(1)
        chrome.append((header, footer))
    return chrome


def _header_chrome_shape(slide, slide_height: int):
    header_line = int(slide_height * HEADER_CHROME_TOP_RATIO)
    for shape in slide.shapes:
        if not shape.has_text_frame or not _is_chrome_shape(shape):
            continue
        text = _shape_text(shape)
        if not text or text.isdigit():
            continue
        if int(shape.top) < header_line:
            return shape
    return None


def _footer_chrome_shape(slide, slide_height: int):
    footer_line = int(slide_height * FOOTER_CHROME_TOP_RATIO)
    for shape in slide.shapes:
        if not shape.has_text_frame or not _is_chrome_shape(shape):
            continue
        text = _shape_text(shape)
        if not text or text.isdigit():
            continue
        if int(shape.top) >= footer_line:
            return shape
    return None


def _set_chrome_shape_text(shape, text: str, color: RGBColor) -> None:
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(CHROME_FONT_PT)
    run.font.color.rgb = color
    run.font.bold = False


def _build_pptx_chrome_map(
    prs: Presentation,
    md_slides: list[str],
    chrome_per_md: list[tuple[str, str]],
) -> list[tuple[str, str]]:
    result: list[tuple[str, str]] = []
    md_index = 0
    pptx_index = 0
    current = chrome_per_md[0] if chrome_per_md else ("", "")

    while pptx_index < len(prs.slides):
        matched = False
        while md_index < len(md_slides):
            current = chrome_per_md[md_index]
            slide_index = _find_pptx_slide_index(
                prs, md_slides[md_index], start_index=pptx_index
            )
            if slide_index == pptx_index:
                result.append(chrome_per_md[md_index])
                md_index += 1
                pptx_index += 1
                matched = True
                break
            if slide_index is None:
                md_index += 1
                continue
            if slide_index > pptx_index:
                break
            md_index += 1
        if not matched:
            result.append(current)
            pptx_index += 1

    return result


def apply_slide_chrome_from_markdown(prs: Presentation, md_path: Path) -> int:
    markdown = md_path.read_text(encoding="utf-8")
    md_slides = split_marp_slides(markdown)
    default_header, default_footer = _parse_marp_frontmatter_chrome(markdown)
    chrome_per_md = _cumulative_slide_chrome(md_slides, default_header, default_footer)
    chrome_map = _build_pptx_chrome_map(prs, md_slides, chrome_per_md)
    slide_height = int(prs.slide_height)
    updated = 0

    for slide_index, (header, footer) in enumerate(chrome_map):
        if slide_index >= len(prs.slides):
            break
        slide = prs.slides[slide_index]
        header_shape = _header_chrome_shape(slide, slide_height)
        footer_shape = _footer_chrome_shape(slide, slide_height)
        if header_shape is not None and _shape_text(header_shape) != header:
            _set_chrome_shape_text(header_shape, header, HEADING_RED)
            updated += 1
        if footer_shape is not None and footer and _shape_text(footer_shape) != footer:
            _set_chrome_shape_text(footer_shape, footer, BODY_BLACK)
            updated += 1

    return updated


def _center_text_frame(text_frame) -> None:
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER


def _left_align_text_frame(text_frame) -> None:
    for paragraph in text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT


def _content_left_margin(slide_width: int) -> int:
    return int(slide_width * CONTENT_LEFT_RATIO)


def _center_shape_horizontally(shape, slide_width: int) -> None:
    shape.left = max(0, int((slide_width - int(shape.width)) // 2))


def _is_horizontally_centered(shape, slide_width: int) -> bool:
    center_x = int(shape.left) + int(shape.width) // 2
    return abs(center_x - slide_width // 2) < slide_width * SECTION_CENTER_TOLERANCE


def _is_module_lead_slide(slide) -> bool:
    if _is_course_title_slide(slide) or _is_lesson_lead_slide(slide):
        return False
    texts: list[str] = []
    has_large_title = False
    for shape in slide.shapes:
        if not shape.has_text_frame or _is_chrome_shape(shape):
            continue
        text = _shape_text(shape)
        if text:
            texts.append(text)
        max_pt = _max_run_font_pt(shape)
        if max_pt is not None and max_pt >= HEADING_FONT_PT - 2:
            has_large_title = True
    if not has_large_title:
        return False
    return any(re.search(r"Module \d+\s*[·•]", text) for text in texts)


def _module_lead_shapes(slide) -> tuple:
    title = module_line = tagline = None
    for shape in slide.shapes:
        if not shape.has_text_frame or _is_chrome_shape(shape):
            continue
        text = _shape_text(shape)
        if re.search(r"Module \d+\s*[·•]", text):
            module_line = shape
        elif (_max_run_font_pt(shape) or 0) >= HEADING_MIN_FONT_PT and title is None:
            title = shape
        elif text and tagline is None and shape is not title and shape is not module_line:
            tagline = shape
    return title, module_line, tagline


def _stack_section_shapes_centered(
    shapes: list,
    slide_height: int,
    slide_width: int,
    *,
    center_text: bool = True,
    gaps: list[int] | None = None,
) -> None:
    if not shapes:
        return
    if gaps is None:
        gaps = [int(Emu(100000)), int(Emu(80000))]

    for shape in shapes:
        _trim_oversized_body_textbox(shape)

    total_height = sum(int(shape.height) for shape in shapes)
    gap_values = gaps[: max(len(shapes) - 1, 0)]
    while len(gap_values) < max(len(shapes) - 1, 0):
        gap_values.append(gaps[-1] if gaps else int(Emu(80000)))
    total_height += sum(gap_values)
    footer_line = int(slide_height * FOOTER_SAFE_RATIO)
    min_start = int(Emu(700000))
    start_y = max(min_start, int((slide_height - total_height) // 2))
    if start_y + total_height > footer_line:
        start_y = max(min_start, footer_line - total_height - int(Emu(50000)))
    cursor = start_y
    for index, shape in enumerate(shapes):
        _center_shape_horizontally(shape, slide_width)
        shape.top = cursor
        if center_text and shape.has_text_frame:
            _center_text_frame(shape.text_frame)
        cursor += int(shape.height)
        if index < len(shapes) - 1:
            cursor += gap_values[index]


def _align_content_slide(slide, slide_width: int) -> None:
    content_left = _content_left_margin(slide_width)

    for shape in _content_layout_shapes(slide):
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            if not _slide_has_non_table_body(slide):
                shape.left = max(0, int((slide_width - int(shape.width)) // 2))
            else:
                shape.left = content_left
            continue
        if not shape.has_text_frame:
            continue
        _left_align_text_frame(shape.text_frame)
        if _is_heading_shape(shape):
            shape.left = content_left
        elif _is_horizontally_centered(shape, slide_width) or int(shape.left) < content_left:
            shape.left = content_left


def _format_module_lead_slide(slide, slide_height: int, slide_width: int) -> None:
    title, module_line, tagline = _module_lead_shapes(slide)
    specs = [
        (title, LEAD_H1_FONT_PT, True, HEADING_RED),
        (module_line, LEAD_H2_FONT_PT, False, HEADING_RED),
        (tagline, MIN_BODY_FONT_PT, False, BODY_BLACK),
    ]
    shapes: list = []
    for shape, target_pt, bold, color in specs:
        if shape is None:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(target_pt)
                run.font.bold = bold
                run.font.color.rgb = color
        shapes.append(shape)

    _stack_section_shapes_centered(shapes, slide_height, slide_width)


def _same_shape(first, second) -> bool:
    if first is None or second is None:
        return first is second
    return first._element is second._element


def _is_up_next_slide(slide) -> bool:
    for shape in slide.shapes:
        if not shape.has_text_frame or _is_chrome_shape(shape):
            continue
        if re.match(r"^Up Next:\s+Module\s+\d+", _shape_text(shape).strip(), re.I):
            return True
    return False


def _up_next_shapes(slide) -> tuple:
    title = subtitle = body = None
    ordered: list = []
    for shape in slide.shapes:
        if not shape.has_text_frame or _is_chrome_shape(shape):
            continue
        text = _shape_text(shape).strip()
        if text:
            ordered.append((int(shape.top), shape, text))
    ordered.sort(key=lambda item: item[0])

    for _, shape, text in ordered:
        if title is None and re.match(r"^Up Next:\s+Module\s+\d+", text, re.I):
            title = shape
        elif (
            subtitle is None
            and not _same_shape(shape, title)
            and (_max_run_font_pt(shape) or 0) >= HEADING_MIN_FONT_PT
            and not text.startswith("Now that")
            and not re.match(r"^End of Module\s+\d+", text, re.I)
        ):
            subtitle = shape
        elif not _same_shape(shape, title) and not _same_shape(shape, subtitle):
            body = shape if body is None else body
    return title, subtitle, body


def _up_next_body_shapes(slide, title, subtitle) -> list:
    bodies: list = []
    for shape in slide.shapes:
        if not shape.has_text_frame or _is_chrome_shape(shape):
            continue
        if _same_shape(shape, title) or _same_shape(shape, subtitle):
            continue
        if _shape_text(shape).strip():
            bodies.append(shape)
    return sorted(bodies, key=lambda shape: (int(shape.top), int(shape.left)))


def _split_up_next_combined_shape(slide, shape, slide_width: int) -> tuple:
    """LibreOffice often exports Up Next slides as one text box — split into title/subtitle/body."""
    paragraphs = [p for p in shape.text_frame.paragraphs if p.text.strip()]
    if len(paragraphs) < 2:
        return None, None, None
    if not re.match(r"^Up Next:\s+Module\s+\d+", paragraphs[0].text.strip(), re.I):
        return None, None, None

    content_left = _content_left_margin(slide_width)
    width = int(slide_width * 0.88)
    line_height = int(Pt(MIN_BODY_FONT_PT * BODY_LINE_SPACING).emu)

    title_text = paragraphs[0].text.strip()
    subtitle_text = ""
    body_texts: list[str] = []
    if len(paragraphs) > 1 and not paragraphs[1].text.strip().startswith("Now that"):
        subtitle_text = paragraphs[1].text.strip()
        body_texts = [p.text.strip() for p in paragraphs[2:]]
    else:
        body_texts = [p.text.strip() for p in paragraphs[1:]]

    element = shape._element
    element.getparent().remove(element)

    def add_textbox(lines: list[str]) -> object | None:
        if not lines:
            return None
        height = len(lines) * line_height + int(Emu(80000))
        box = slide.shapes.add_textbox(content_left, 0, width, height)
        frame = box.text_frame
        frame.word_wrap = True
        for index, line in enumerate(lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = line
        return box

    title = add_textbox([title_text])
    subtitle = add_textbox([subtitle_text]) if subtitle_text else None
    body = add_textbox(body_texts) if body_texts else None
    return title, subtitle, body


def _format_up_next_slide(slide, slide_height: int, slide_width: int) -> None:
    title, subtitle, _body = _up_next_shapes(slide)
    body_shapes = _up_next_body_shapes(slide, title, subtitle)

    if title is not None and subtitle is None and not body_shapes:
        split = _split_up_next_combined_shape(slide, title, slide_width)
        if split[0] is not None:
            title, subtitle, body = split
            body_shapes = [body] if body is not None else []
        else:
            body = None
    else:
        if len(body_shapes) > 1:
            _merge_body_shape_group(slide, body_shapes, slide_width=slide_width)
            body_shapes = _up_next_body_shapes(slide, title, subtitle)
        body = body_shapes[0] if body_shapes else None
    specs = [
        (title, LEAD_H1_FONT_PT, True, HEADING_RED),
        (subtitle, LEAD_H2_FONT_PT, False, HEADING_RED),
    ]
    shapes: list = []
    for shape, target_pt, bold, color in specs:
        if shape is None:
            continue
        for paragraph in shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(target_pt)
                run.font.bold = bold
                run.font.color.rgb = color
        shapes.append(shape)

    if body is not None:
        _trim_oversized_body_textbox(body)
        max_width = int(slide_width * 0.92)
        if int(body.width) > max_width:
            body.width = max_width
        for paragraph in body.text_frame.paragraphs:
            text = paragraph.text.strip()
            is_tagline = bool(re.match(r"^End of Module\s+\d+", text, re.I))
            _apply_body_paragraph_spacing(paragraph)
            paragraph.alignment = PP_ALIGN.CENTER
            if is_tagline:
                paragraph.space_before = Pt(10)
            for run in paragraph.runs:
                run.font.size = Pt(MIN_BODY_FONT_PT)
                run.font.color.rgb = BODY_BLACK
                run.font.bold = False
                run.font.italic = is_tagline
        shapes.append(body)

    _stack_section_shapes_centered(
        shapes,
        slide_height,
        slide_width,
        gaps=[int(Emu(80000)), int(Emu(140000)), int(Emu(80000))],
    )

    footer_line = int(slide_height * FOOTER_SAFE_RATIO)
    body_pt = MIN_BODY_FONT_PT
    for _ in range(6):
        if body is None:
            break
        max_bottom = max(_shape_bounds(shape)[3] for shape in shapes)
        if max_bottom <= footer_line:
            break
        body_pt = max(13, body_pt - 1)
        for paragraph in body.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(body_pt)
        _fit_textbox_height(body)
        _stack_section_shapes_centered(
            shapes,
            slide_height,
            slide_width,
            gaps=[int(Emu(60000)), int(Emu(100000)), int(Emu(60000))],
        )

    max_bottom = max(_shape_bounds(shape)[3] for shape in shapes)
    if max_bottom > footer_line:
        shift = max_bottom - footer_line + int(Emu(20000))
        for shape in shapes:
            shape.top = max(int(Emu(500000)), int(shape.top) - shift)


def _is_course_title_slide(slide) -> bool:
    title = subtitle = False
    for shape in slide.shapes:
        if not shape.has_text_frame or _is_chrome_shape(shape):
            continue
        text = _shape_text(shape)
        if text == "Cursor Training Program":
            title = True
        if "AI-Assisted Development" in text:
            subtitle = True
    return title and subtitle


def _is_course_agenda_slide(slide) -> bool:
    heading = _pptx_slide_heading(slide)
    if not heading:
        return False
    return heading.startswith("course agenda") or heading.startswith("day 1 —") or heading.startswith("day 2 —")


def _is_lesson_lead_slide(slide) -> bool:
    for shape in slide.shapes:
        if not shape.has_text_frame or _is_chrome_shape(shape):
            continue
        if re.match(r"^Lesson \d+\.\d+$", _shape_text(shape).strip()):
            return True
    return False


def _is_lead_slide(slide) -> bool:
    if _is_course_title_slide(slide):
        return True
    if _is_lesson_lead_slide(slide):
        return True
    texts: list[str] = []
    has_large_title = False
    for shape in slide.shapes:
        if not shape.has_text_frame or _is_chrome_shape(shape):
            continue
        text = _shape_text(shape)
        if text:
            texts.append(text)
        max_pt = _max_run_font_pt(shape)
        if max_pt is not None and max_pt >= HEADING_FONT_PT - 2:
            has_large_title = True
    if not has_large_title:
        return False
    return any(re.search(r"Module \d+\s*[·•]", text) for text in texts)


def _normalize_slide_headings(slide) -> None:
    if _is_course_title_slide(slide):
        return
    if _is_up_next_slide(slide):
        return
    if _is_lesson_lead_slide(slide):
        lesson_shape = None
        topic_shape = None
        for shape in slide.shapes:
            if not shape.has_text_frame or _is_chrome_shape(shape):
                continue
            text = _shape_text(shape).strip()
            if re.match(r"^Lesson \d+\.\d+$", text):
                lesson_shape = shape
            elif (
                lesson_shape is not None
                and shape is not lesson_shape
                and (_max_run_font_pt(shape) or 0) >= HEADING_MIN_FONT_PT
                and topic_shape is None
            ):
                topic_shape = shape
        if lesson_shape is not None:
            for paragraph in lesson_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(HEADING_FONT_PT)
                    run.font.bold = True
                    run.font.color.rgb = HEADING_RED
        if topic_shape is not None:
            for paragraph in topic_shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(HEADING_FONT_PT)
                    run.font.bold = False
                    run.font.color.rgb = HEADING_RED
        return
    is_lead = _is_lead_slide(slide)
    for shape in slide.shapes:
        if not _is_heading_shape(shape):
            continue
        text = _shape_text(shape)
        if is_lead and re.search(r"Module \d+\s*[·•]", text):
            target_pt = LEAD_H2_FONT_PT
            bold = False
        elif is_lead:
            target_pt = LEAD_H1_FONT_PT
            bold = True
        else:
            target_pt = HEADING_FONT_PT
            bold = True

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(target_pt)
                run.font.bold = bold
                run.font.color.rgb = HEADING_RED


def _course_title_shapes(slide) -> tuple:
    title = subtitle = tagline = None
    for shape in slide.shapes:
        if not shape.has_text_frame or _is_chrome_shape(shape):
            continue
        text = _shape_text(shape)
        if text == "Cursor Training Program":
            title = shape
        elif "AI-Assisted Development" in text:
            subtitle = shape
        elif "Springpeople ·" in text or "2-day instructor" in text.lower():
            tagline = shape
    return title, subtitle, tagline


def _format_course_title_slide(slide, slide_height: int, slide_width: int) -> None:
    title, subtitle, tagline = _course_title_shapes(slide)
    if title is None:
        return

    specs = [
        (title, LEAD_H1_FONT_PT, True, HEADING_RED),
        (subtitle, LEAD_H2_FONT_PT, False, HEADING_RED),
        (tagline, MIN_BODY_FONT_PT, False, BODY_BLACK),
    ]
    stacked: list = []
    for shape, target_pt, bold, color in specs:
        if shape is None:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(target_pt)
                run.font.bold = bold
                run.font.color.rgb = color
        stacked.append(shape)

    _stack_section_shapes_centered(stacked, slide_height, slide_width)


def _course_agenda_body_shape(slide):
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame or _is_chrome_shape(shape) or _is_heading_shape(shape):
            continue
        text = _shape_text(shape)
        if not text:
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            continue
        candidates.append(shape)
    if not candidates:
        return None
    return max(candidates, key=lambda shape: int(shape.top))


def _shrink_table_to_height(shape, max_height: int) -> None:
    table = shape.table
    if max_height <= 0 or int(shape.height) <= max_height:
        return
    scale = max_height / int(shape.height)
    for row in table.rows:
        row.height = max(int(Emu(70000)), int(int(row.height) * scale))
    font_pt = max(11, TABLE_FONT_PT * scale)
    line_pt = max(14, TABLE_ROW_HEIGHT_PT * scale)
    for row_idx in range(len(table.rows)):
        for col_idx in range(len(table.columns)):
            cell = table.cell(row_idx, col_idx)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.line_spacing = Pt(line_pt)
                for run in paragraph.runs:
                    run.font.size = Pt(font_pt)
    shape.height = sum(int(row.height) for row in table.rows)


def _format_course_agenda_slide(slide, slide_height: int, slide_width: int) -> None:
    heading = None
    table = None
    for shape in slide.shapes:
        if _is_chrome_shape(shape):
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table = shape
        elif shape.has_text_frame and _is_heading_shape(shape):
            heading = shape

    body = _course_agenda_body_shape(slide)
    heading_top = int(Emu(1200000))
    footer_line = int(slide_height * FOOTER_SAFE_RATIO)

    if heading is not None:
        heading.top = heading_top
        _trim_oversized_body_textbox(heading)
        for paragraph in heading.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(HEADING_FONT_PT)
                run.font.bold = True
                run.font.color.rgb = HEADING_RED
        cursor = int(heading.top) + int(heading.height) + HEADING_BODY_GAP
    else:
        cursor = heading_top

    if body is not None:
        _trim_oversized_body_textbox(body)
        for paragraph in body.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(MIN_BODY_FONT_PT)
                run.font.color.rgb = BODY_BLACK

    body_height = int(body.height) if body is not None else 0
    body_gap = int(CONTENT_GAP) if body is not None else 0
    max_table_bottom = footer_line - body_height - body_gap

    if table is not None:
        table.top = cursor
        table.left = max(0, int((slide_width - int(table.width)) // 2))
        available = max_table_bottom - int(table.top)
        _shrink_table_to_height(table, available)
        cursor = int(table.top) + int(table.height) + CONTENT_GAP

    if body is not None:
        body.top = min(cursor, footer_line - body_height)
        for paragraph in body.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = "Total:" in paragraph.text


def _format_lesson_lead_slide(slide, slide_height: int, slide_width: int) -> None:
    ordered: list = []
    for shape in slide.shapes:
        if not shape.has_text_frame or _is_chrome_shape(shape):
            continue
        text = _shape_text(shape).strip()
        if re.match(r"^Lesson \d+\.\d+$", text):
            ordered.append((0, shape))
        elif (_max_run_font_pt(shape) or 0) >= HEADING_MIN_FONT_PT:
            ordered.append((1, shape))
        elif text:
            ordered.append((2, shape))
    ordered.sort(key=lambda item: item[0])
    shapes = [shape for _, shape in ordered]
    if not shapes:
        return

    for shape in shapes:
        _trim_oversized_body_textbox(shape)
        if (_max_run_font_pt(shape) or 0) < HEADING_MIN_FONT_PT:
            _normalize_body_fonts(shape)

    _stack_section_shapes_centered(shapes, slide_height, slide_width)


def _format_course_intro_slides(prs: Presentation) -> None:
    slide_height = int(prs.slide_height)
    slide_width = int(prs.slide_width)
    for slide in prs.slides:
        if _is_course_title_slide(slide):
            _format_course_title_slide(slide, slide_height, slide_width)
        elif _is_course_agenda_slide(slide):
            _format_course_agenda_slide(slide, slide_height, slide_width)


def _normalize_slide_formatting(slide) -> None:
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            _normalize_table_fonts(shape)
            continue
        if shape.has_text_frame:
            if _is_code_shape(shape):
                _normalize_code_block_fonts(shape)
            elif not _is_chrome_shape(shape):
                if not _is_heading_shape(shape):
                    _normalize_body_fonts(shape)
    _normalize_slide_headings(slide)


def _normalize_presentation_formatting(prs: Presentation) -> None:
    for slide in prs.slides:
        _normalize_slide_formatting(slide)


def _normalize_text(text: str) -> str:
    text = re.sub(r"`([^`]*)`", r"\1", text)
    text = re.sub(r"\*\*", "", text)
    text = text.replace("→", "->")
    return " ".join(text.split()).lower()


def _markdown_slide_heading(slide_md: str) -> str | None:
    for line in slide_md.splitlines():
        if line.startswith("## "):
            return _normalize_text(line[3:])
    for line in slide_md.splitlines():
        if line.startswith("# "):
            return _normalize_text(line[2:])
    return None


def _pptx_slide_heading(slide) -> str | None:
    best_text: str | None = None
    best_size = 0.0
    for shape in slide.shapes:
        if not shape.has_text_frame or _is_chrome_shape(shape):
            continue
        text = _shape_text(shape).split("\n", 1)[0].strip()
        if not text:
            continue
        size = _max_run_font_pt(shape) or 0.0
        if size >= HEADING_MIN_FONT_PT and size >= best_size:
            best_size = size
            best_text = _normalize_text(text)
    return best_text


def _find_pptx_slide_index(
    prs: Presentation,
    slide_md: str,
    *,
    start_index: int = 0,
) -> int | None:
    target = _markdown_slide_heading(slide_md)
    if not target:
        return None

    for index in range(start_index, len(prs.slides)):
        heading = _pptx_slide_heading(prs.slides[index])
        if heading == target:
            return index

    best_index: int | None = None
    best_score = 0
    for index in range(start_index, len(prs.slides)):
        heading = _pptx_slide_heading(prs.slides[index])
        if not heading:
            continue
        if target in heading or heading in target:
            score = min(len(target), len(heading))
            if score > best_score:
                best_score = score
                best_index = index
    return best_index


def _cell_match_score(shape_text: str, cell_text: str) -> int:
    left = _normalize_text(shape_text)
    right = _normalize_text(cell_text)
    if not left or not right:
        return 0
    if left == right:
        return 10_000
    if left in right and len(left) < len(right) * 0.95:
        return len(left)
    if right in left and len(left) <= len(right) * 1.15:
        return len(right)
    return 0


def _table_cell_texts(table_data: list[list[str]]) -> list[str]:
    return [cell for row in table_data for cell in row if cell.strip()]


def _shape_matches_table(shape, table_data: list[list[str]], min_score: int = 10) -> bool:
    text = _shape_text(shape)
    if not text:
        return False
    score = max(
        (_cell_match_score(text, cell) for cell in _table_cell_texts(table_data)),
        default=0,
    )
    if score >= min_score:
        return True
    if len(text.strip()) <= 12:
        return score >= max(8, len(text.strip()))
    return False


def _content_text_shapes(slide) -> list:
    return [
        shape
        for shape in slide.shapes
        if shape.has_text_frame and _shape_text(shape) and not _is_chrome_shape(shape)
    ]


def _shape_bounds(shape) -> tuple[int, int, int, int]:
    return (
        int(shape.left),
        int(shape.top),
        int(shape.left + shape.width),
        int(shape.top + shape.height),
    )


def _bounds_for_shapes(shapes) -> tuple[int, int, int, int]:
    left = min(_shape_bounds(shape)[0] for shape in shapes)
    top = min(_shape_bounds(shape)[1] for shape in shapes)
    right = max(_shape_bounds(shape)[2] for shape in shapes)
    bottom = max(_shape_bounds(shape)[3] for shape in shapes)
    return left, top, right - left, bottom - top


def _shape_in_bounds(shape, bounds: tuple[int, int, int, int], padding: int = 80000) -> bool:
    left, top, width, height = bounds
    right = left + width
    bottom = top + height
    sl, st, sr, sb = _shape_bounds(shape)
    return (
        sl >= left - padding
        and sr <= right + padding
        and st >= top - padding
        and sb <= bottom + padding
    )


def _match_table_shapes(available, assigned_ids: set[int], table_data: list[list[str]]):
    matched = [
        shape
        for shape in available
        if id(shape) not in assigned_ids and _shape_matches_table(shape, table_data)
    ]
    if not matched:
        matched = [
            shape
            for shape in available
            if id(shape) not in assigned_ids
            and max(
                (
                    _cell_match_score(_shape_text(shape), cell)
                    for cell in _table_cell_texts(table_data)
                ),
                default=0,
            )
            >= 8
        ]
    if not matched:
        return []

    bounds = _bounds_for_shapes(matched)
    expanded = list(matched)
    for shape in available:
        if id(shape) in assigned_ids or shape in expanded:
            continue
        if _shape_in_bounds(shape, bounds):
            expanded.append(shape)
    return expanded


def _remove_shapes(shapes) -> None:
    for shape in shapes:
        element = shape._element
        element.getparent().remove(element)


LIST_MARKER_RE = re.compile(r"^[-•*–]\s+")
BULLET_CHAR = "•"
# Standard PowerPoint first-level hanging bullet (marL + indent positions bullet vs text).
BULLET_LEVEL0_MAR_L = Emu(342900)  # ~0.375 in — text start
BULLET_HANGING = Emu(-342900)  # bullet at text-box left; full marL width before text
BULLET_LEVEL_STEP = Emu(400050)  # extra marL per nested list level
_BULLET_TAG_NAMES = (
    qn("a:buChar"),
    qn("a:buNone"),
    qn("a:buAutoNum"),
    qn("a:buFont"),
    qn("a:buBlip"),
)


def _paragraph_text(paragraph) -> str:
    return "".join(run.text for run in paragraph.runs if run.text).strip()


def _looks_like_list_item(text: str) -> bool:
    return bool(LIST_MARKER_RE.match(text.strip()))


def _clear_paragraph_bullet(pPr) -> None:
    for child in list(pPr):
        if child.tag in _BULLET_TAG_NAMES:
            pPr.remove(child)


def _set_paragraph_indents(pPr, mar_l: int, indent: int) -> None:
    for tag, val in ((qn("a:marL"), mar_l), (qn("a:indent"), indent)):
        el = pPr.find(tag)
        if el is None:
            local = "marL" if tag.endswith("marL") else "indent"
            el = OxmlElement(f"a:{local}")
            pPr.append(el)
        el.set("val", str(val))


def _apply_body_paragraph_spacing(paragraph) -> None:
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = BODY_LINE_SPACING


def _enable_paragraph_bullet(paragraph, *, char: str = BULLET_CHAR) -> None:
    pPr = paragraph._p.get_or_add_pPr()
    _clear_paragraph_bullet(pPr)
    bu_char = OxmlElement("a:buChar")
    bu_char.set("char", char)
    pPr.insert(0, bu_char)
    level = max(0, int(paragraph.level or 0))
    mar_l = int(BULLET_LEVEL0_MAR_L + BULLET_LEVEL_STEP * level)
    _set_paragraph_indents(pPr, mar_l, int(BULLET_HANGING))


def _prepend_bullet_text_gap(paragraph) -> None:
    """Leading space after buChar — reliable gap when theme lstStyle ignores indents."""
    for run in paragraph.runs:
        if not run.text:
            continue
        if not run.text.startswith((" ", "\u00a0", "\t")):
            run.text = " " + run.text
        break


def _strip_list_marker_from_paragraph(paragraph) -> None:
    for run in paragraph.runs:
        if not run.text:
            continue
        updated = LIST_MARKER_RE.sub("", run.text, count=1)
        if updated != run.text:
            run.text = updated
        break


def _apply_body_list_formatting(shape) -> None:
    if not shape.has_text_frame or _is_chrome_shape(shape) or _is_heading_shape(shape):
        return
    if _is_code_shape(shape):
        return
    for paragraph in shape.text_frame.paragraphs:
        text = _paragraph_text(paragraph)
        if text:
            _apply_body_paragraph_spacing(paragraph)
        if not text:
            continue
        is_list = _looks_like_list_item(text)
        if is_list:
            paragraph.level = 0
            _strip_list_marker_from_paragraph(paragraph)
            _enable_paragraph_bullet(paragraph)
            _prepend_bullet_text_gap(paragraph)
        elif paragraph.level > 0:
            _strip_list_marker_from_paragraph(paragraph)
            _enable_paragraph_bullet(paragraph)
            _prepend_bullet_text_gap(paragraph)


def _copy_run_font(source_run, target_run) -> None:
    if source_run.font.name:
        target_run.font.name = source_run.font.name
    if source_run.font.size:
        target_run.font.size = source_run.font.size
    if source_run.font.bold is not None:
        target_run.font.bold = source_run.font.bold
    if source_run.font.italic is not None:
        target_run.font.italic = source_run.font.italic
    try:
        if source_run.font.color.rgb:
            target_run.font.color.rgb = source_run.font.color.rgb
    except Exception:
        pass


def _append_shape_runs_to_paragraph(paragraph, shape) -> None:
    for src_index, src_paragraph in enumerate(shape.text_frame.paragraphs):
        for src_run in src_paragraph.runs:
            if not src_run.text:
                continue
            run = paragraph.add_run()
            run.text = src_run.text
            _copy_run_font(src_run, run)
        if src_index < len(shape.text_frame.paragraphs) - 1:
            run = paragraph.add_run()
            run.text = "\n"


def _mergeable_body_shapes(slide) -> list:
    shapes = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            continue
        if not shape.has_text_frame:
            continue
        if _is_chrome_shape(shape) or _is_heading_shape(shape):
            continue
        if _is_code_shape(shape) and not _is_inline_monospace_shape(shape):
            continue
        if not _shape_text(shape):
            continue
        shapes.append(shape)
    return shapes


def _body_merge_blockers(slide) -> list:
    blockers = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            blockers.append(shape)
        elif _is_code_shape(shape) and not _is_inline_monospace_shape(shape):
            blockers.append(shape)
    return blockers


def _has_blocker_between(blockers, lower_y: int, upper_y: int) -> bool:
    if upper_y <= lower_y:
        return False
    for blocker in blockers:
        blocker_top = int(blocker.top)
        if lower_y < blocker_top < upper_y:
            return True
    return False


def _group_body_shapes_for_merge(shapes: list, blockers: list) -> list[list]:
    ordered = sorted(shapes, key=lambda shape: (int(shape.top), int(shape.left)))
    groups: list[list] = []
    current: list = []
    previous_bottom = 0

    for shape in ordered:
        shape_top = int(shape.top)
        if current and _has_blocker_between(blockers, previous_bottom, shape_top):
            groups.append(current)
            current = [shape]
        else:
            current.append(shape)
        previous_bottom = max(previous_bottom, int(shape.top) + int(shape.height))

    if current:
        groups.append(current)
    return groups


def _absorb_inline_monospace_shapes(slide) -> None:
    for shape in slide.shapes:
        if _is_inline_monospace_shape(shape):
            _convert_shape_to_body_text(shape)


def _same_text_line(first, second, tolerance: int | None = None) -> bool:
    if tolerance is None:
        tolerance = int(Pt(MIN_BODY_FONT_PT * BODY_LINE_SPACING * 0.65).emu)
    return abs(int(first.top) - int(second.top)) <= tolerance


def _replace_shape_text(shape, text: str) -> None:
    if not shape.has_text_frame:
        return
    frame = shape.text_frame
    frame.clear()
    paragraph = frame.paragraphs[0]
    _apply_body_paragraph_spacing(paragraph)
    run = paragraph.add_run()
    run.text = text
    run.font.name = BODY_FONT
    run.font.size = Pt(MIN_BODY_FONT_PT)
    run.font.color.rgb = BODY_BLACK


def _coalesce_same_line_body_fragments(slide) -> None:
    """Merge LibreOffice inline-backtick fragments that land on one horizontal line."""
    shapes = _mergeable_body_shapes(slide)
    if len(shapes) < 2:
        return

    ordered = sorted(shapes, key=lambda shape: (int(shape.top), int(shape.left)))
    groups: list[list] = []
    current = [ordered[0]]
    for shape in ordered[1:]:
        previous = current[-1]
        prev_right = int(previous.left) + int(previous.width)
        same_line = _same_text_line(previous, shape)
        adjacent = int(shape.left) <= prev_right + int(Emu(450000))
        if same_line and adjacent:
            current.append(shape)
        else:
            groups.append(current)
            current = [shape]
    groups.append(current)

    for group in groups:
        if len(group) < 2:
            continue
        group.sort(key=lambda shape: int(shape.left))
        combined_parts: list[str] = []
        for index, shape in enumerate(group):
            text = _shape_text(shape)
            if not text:
                continue
            if index > 0 and combined_parts:
                previous = combined_parts[-1]
                if not previous.endswith((" ", "\n")) and not text.startswith(
                    (",", ".", ";", ":", ")", "—", "·", "?", "!")
                ):
                    combined_parts.append(" ")
            combined_parts.append(text)
        combined = "".join(combined_parts).strip()
        if not combined:
            continue
        _replace_shape_text(group[0], combined)
        _remove_shapes(group[1:])


def _should_skip_body_merge(slide) -> bool:
    return (
        _is_course_title_slide(slide)
        or _is_course_agenda_slide(slide)
        or _is_lesson_lead_slide(slide)
        or _is_module_lead_slide(slide)
        or _is_up_next_slide(slide)
    )


def _sort_shapes_by_markdown_order(shapes: list, slide_md: str | None) -> list:
    if not slide_md:
        return sorted(shapes, key=lambda shape: (int(shape.top), int(shape.left)))

    rank: dict[int, int] = {}
    for index, (kind, text) in enumerate(_markdown_content_sequence(slide_md)):
        for shape in shapes:
            if _sequence_item_match_score(shape, kind, text) <= 0:
                continue
            key = id(shape._element)
            rank[key] = min(rank.get(key, index), index)
    return sorted(
        shapes,
        key=lambda shape: (rank.get(id(shape._element), 9999), int(shape.top), int(shape.left)),
    )


def _merge_body_shape_group(
    slide,
    shapes: list,
    *,
    slide_width: int,
    slide_md: str | None = None,
) -> None:
    if len(shapes) < 2:
        return

    ordered = _sort_shapes_by_markdown_order(shapes, slide_md)
    left = min(int(shape.left) for shape in ordered)
    top = min(int(shape.top) for shape in ordered)
    content_left = _content_left_margin(slide_width)
    width = max(int(slide_width * 0.88), max(int(s.left) + int(s.width) for s in ordered) - left)
    line_height = int(Pt(MIN_BODY_FONT_PT * BODY_LINE_SPACING).emu)
    height = len(ordered) * line_height + int(Emu(40000))

    textbox = slide.shapes.add_textbox(content_left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.margin_left = Emu(0)
    frame.margin_right = Emu(0)
    frame.margin_top = Emu(0)
    frame.margin_bottom = Emu(0)

    bullet_threshold = content_left + int(Emu(350000))
    for index, shape in enumerate(ordered):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        _apply_body_paragraph_spacing(paragraph)
        shape_text = _shape_text(shape)
        is_bullet = int(shape.left) >= bullet_threshold or _looks_like_list_item(shape_text)
        if is_bullet:
            paragraph.level = 0
        _append_shape_runs_to_paragraph(paragraph, shape)
        if is_bullet:
            _strip_list_marker_from_paragraph(paragraph)
            _enable_paragraph_bullet(paragraph)
            _prepend_bullet_text_gap(paragraph)
        elif paragraph.level > 0:
            _strip_list_marker_from_paragraph(paragraph)
            _enable_paragraph_bullet(paragraph)
            _prepend_bullet_text_gap(paragraph)

    _apply_body_list_formatting(textbox)
    _remove_shapes(ordered)


def _merge_body_textboxes_on_slide(
    slide,
    *,
    slide_width: int,
    slide_md: str | None = None,
) -> int:
    if _should_skip_body_merge(slide):
        return 0

    shapes = _mergeable_body_shapes(slide)
    if len(shapes) < 2:
        return 0

    blockers = _body_merge_blockers(slide)
    groups = _group_body_shapes_for_merge(shapes, blockers)
    merged = 0
    for group in groups:
        if len(group) < 2:
            continue
        _merge_body_shape_group(slide, group, slide_width=slide_width, slide_md=slide_md)
        merged += 1
    return merged


def merge_body_textboxes(prs: Presentation, md_slides: list[str] | None = None) -> int:
    slide_width = int(prs.slide_width)
    merged = 0
    for idx, slide in enumerate(prs.slides):
        slide_md = md_slides[idx] if md_slides and idx < len(md_slides) else None
        merged += _merge_body_textboxes_on_slide(
            slide, slide_width=slide_width, slide_md=slide_md
        )
    return merged


def _set_cell_fill(cell, color: RGBColor) -> None:
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_cell_border(cell, color: str = TABLE_BORDER, width: str = "12700") -> None:
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    for edge in ("lnL", "lnR", "lnT", "lnB"):
        ln = OxmlElement(f"a:{edge}")
        ln.set("w", width)
        ln.set("cap", "flat")
        ln.set("cmpd", "sng")
        ln.set("algn", "ctr")

        solid_fill = OxmlElement("a:solidFill")
        srgb = OxmlElement("a:srgbClr")
        srgb.set("val", color)
        solid_fill.append(srgb)
        ln.append(solid_fill)

        dash = OxmlElement("a:prstDash")
        dash.set("val", "solid")
        ln.append(dash)

        tc_pr.append(ln)


def _clean_table_cell_text(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"\*(.+?)\*", r"\1", text)
    return text.strip()


def _format_table_cell(
    cell,
    text: str,
    *,
    row_idx: int,
    header_row: bool,
) -> None:
    frame = cell.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = TABLE_CELL_INSET
    frame.margin_right = TABLE_CELL_INSET
    frame.margin_top = Emu(15000)
    frame.margin_bottom = Emu(15000)
    paragraph = frame.paragraphs[0]
    paragraph.line_spacing = TABLE_LINE_SPACING
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = _clean_table_cell_text(text)
    run.font.size = Pt(TABLE_FONT_PT)
    run.font.bold = header_row
    run.font.color.rgb = HEADING_RED if header_row else BODY_BLACK

    if header_row:
        fill_color = TABLE_HEADER_BG
    elif row_idx % 2 == 1:
        fill_color = TABLE_ROW_BG
    else:
        fill_color = TABLE_ROW_ALT_BG

    _set_cell_fill(cell, fill_color)
    _set_cell_border(cell)


def _table_dimensions(
    table_data: list[list[str]],
    *,
    max_width: int | None = None,
) -> tuple[list[int], list[int]]:
    rows = len(table_data)
    cols = max(len(row) for row in table_data)
    char_width = int(Pt(TABLE_CHAR_WIDTH_PT).emu)
    line_height = int(Pt(TABLE_ROW_HEIGHT_PT).emu)
    inset = int(TABLE_CELL_INSET) * 2
    buffer = int(Pt(TABLE_SIZE_BUFFER_PT).emu)

    col_max_chars: list[int] = []
    for col_idx in range(cols):
        max_chars = 0
        for row in table_data:
            value = _clean_table_cell_text(row[col_idx] if col_idx < len(row) else "")
            max_chars = max(max_chars, len(value))
        col_max_chars.append(max_chars)

    col_widths: list[int] = []
    for max_chars in col_max_chars:
        effective_chars = min(max_chars, TABLE_MAX_COL_CHARS)
        width = max(int(TABLE_MIN_COL_WIDTH), effective_chars * char_width + inset + buffer)
        col_widths.append(width)

    if max_width is not None:
        while sum(col_widths) > max_width:
            widest = max(range(cols), key=lambda idx: col_widths[idx])
            reduced = col_widths[widest] - char_width
            if reduced <= int(TABLE_MIN_COL_WIDTH):
                break
            col_widths[widest] = max(int(TABLE_MIN_COL_WIDTH), reduced)

    row_heights: list[int] = []
    for row in table_data:
        max_lines = 1
        for col_idx in range(cols):
            value = _clean_table_cell_text(row[col_idx] if col_idx < len(row) else "")
            usable = max(char_width, col_widths[col_idx] - inset)
            chars_per_line = max(1, usable // char_width)
            line_count = max(1, (len(value) + chars_per_line - 1) // chars_per_line)
            max_lines = max(max_lines, line_count)
        row_heights.append(max_lines * line_height + inset + buffer)

    return col_widths, row_heights


def _insert_table(
    slide,
    table_data: list[list[str]],
    bounds: tuple[int, int, int, int],
    *,
    slide_width: int | None = None,
) -> None:
    rows = len(table_data)
    cols = max(len(row) for row in table_data)
    left, top, _, _ = bounds
    max_width = int(slide_width * 0.92) if slide_width else None
    col_widths, row_heights = _table_dimensions(table_data, max_width=max_width)
    width = sum(col_widths)
    height = sum(row_heights)

    graphic = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = graphic.table

    for col_idx, col_width in enumerate(col_widths):
        table.columns[col_idx].width = col_width
    for row_idx, row_height in enumerate(row_heights):
        table.rows[row_idx].height = row_height

    for row_idx, row in enumerate(table_data):
        for col_idx in range(cols):
            value = row[col_idx] if col_idx < len(row) else ""
            _format_table_cell(
                table.cell(row_idx, col_idx),
                value,
                row_idx=row_idx,
                header_row=row_idx == 0,
            )


def _rebuild_slide_tables(slide, tables: list[list[list[str]]], *, slide_width: int | None = None) -> int:
    rebuilt = 0
    assigned_ids: set[int] = set()
    available = _content_text_shapes(slide)

    for table_data in tables:
        matched = _match_table_shapes(available, assigned_ids, table_data)
        if not matched:
            continue

        for shape in matched:
            assigned_ids.add(id(shape))

        bounds = _bounds_for_shapes(matched)
        _remove_shapes(matched)
        _insert_table(slide, table_data, bounds, slide_width=slide_width)
        rebuilt += 1

    return rebuilt


def _shape_uses_monospace(shape) -> bool:
    if not shape.has_text_frame:
        return False
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            name = (run.font.name or "").lower()
            if "consolas" in name or "courier" in name or "mono" in name:
                return True
    return False


def _shape_has_code_fill(shape) -> bool:
    try:
        if shape.fill.type and str(shape.fill.fore_color.rgb).upper() == "F4F4F4":
            return True
    except Exception:
        pass
    return False


def _looks_like_fenced_code(text: str) -> bool:
    stripped = text.strip()
    if not stripped or re.match(r"^Step \d+:", stripped):
        return False
    indicators = (
        "┌",
        "│",
        "└",
        "import ",
        "def ",
        "class ",
        "$ ",
        "curl ",
        "git ",
        "npm ",
        "const ",
        "function ",
        "SELECT ",
        "#!/",
        "POST ",
        "GET ",
        "api.",
        "http",
        "→",
        "->",
        "{",
        "};",
    )
    if any(token in text for token in indicators):
        return True
    lines = [line for line in text.splitlines() if line.strip()]
    if len(lines) >= 3 and sum(
        1 for line in lines if line.startswith((" ", "\t", "•", "·", "-", "*"))
    ) >= 2:
        return True
    if len(lines) == 1 and (stripped.startswith(('"', "'", "`")) or ":" in stripped):
        return True
    return False


def _is_inline_monospace_shape(shape) -> bool:
    if not shape.has_text_frame or _is_chrome_shape(shape):
        return False
    if not _shape_uses_monospace(shape) or _shape_has_code_fill(shape):
        return False
    text = _shape_text(shape)
    if not text or "\n" in text:
        return False
    return len(text) <= 100


def _is_code_shape(shape) -> bool:
    if _is_chrome_shape(shape):
        return False
    if _is_inline_monospace_shape(shape):
        return False
    if _shape_has_code_fill(shape):
        return True
    if _shape_uses_monospace(shape):
        text = _shape_text(shape)
        if not text:
            return False
        if "\n" not in text:
            return False
        return _looks_like_fenced_code(text)
    return False


def _shape_matches_code_content(
    shape_text: str,
    code_text: str,
    *,
    min_ratio: float = 0.5,
) -> bool:
    code_lines = [line.rstrip() for line in code_text.splitlines() if line.strip()]
    if not code_lines:
        return False
    shape_lines = [line.rstrip() for line in shape_text.splitlines() if line.strip()]
    if not shape_lines:
        return False
    matched = sum(
        1
        for code_line in code_lines
        if any(_text_matches_code_line(shape_line, code_line) for shape_line in shape_lines)
    )
    required = max(1, round(len(code_lines) * min_ratio))
    return matched >= required


def _find_fallback_code_shapes(
    available,
    assigned_ids: set[int],
    code_text: str,
) -> list:
    matched = [
        shape
        for shape in available
        if id(shape) not in assigned_ids
        and not _is_heading_shape(shape)
        and _shape_matches_code_content(_shape_text(shape), code_text, min_ratio=0.6)
    ]
    if matched:
        return matched

    code_lines = [line.rstrip() for line in code_text.splitlines() if line.strip()]
    if len(code_lines) == 1:
        target = _normalize_text(code_lines[0])
        for shape in available:
            if id(shape) in assigned_ids or _is_heading_shape(shape):
                continue
            shape_text = _normalize_text(_shape_text(shape))
            if shape_text == target or target in shape_text:
                matched.append(shape)
    return matched


def _estimate_code_insert_top(slide, slide_height: int) -> int:
    shapes = _content_layout_shapes(slide)
    headings = [shape for shape in shapes if _is_heading_shape(shape)]
    if headings:
        primary = min(headings, key=lambda shape: int(shape.top))
        _fit_textbox_height(primary)
        return int(primary.top) + int(primary.height) + HEADING_BODY_GAP
    return int(HEADING_TOP_CONTENT) + int(HEADING_BODY_GAP)


def _text_matches_code_line(shape_text: str, code_line: str) -> bool:
    left = shape_text.rstrip()
    right = code_line.rstrip()
    if left == right:
        return True
    return _normalize_text(left) == _normalize_text(right)


def _code_block_match_score(shape_text: str, code_text: str) -> int:
    if not shape_text.strip() or not code_text.strip():
        return 0
    if _normalize_text(shape_text) == _normalize_text(code_text):
        return 10_000

    shape_lines = [line.rstrip() for line in shape_text.splitlines()]
    code_lines = [line.rstrip() for line in code_text.splitlines() if line.strip()]
    if not shape_lines or not code_lines:
        return 0

    score = 0
    for code_line in code_lines:
        for shape_line in shape_lines:
            if _text_matches_code_line(shape_line, code_line):
                score += max(len(code_line), 8)
                break
    return score


def _match_code_block_shapes(
    available,
    assigned_ids: set[int],
    code_text: str,
) -> list:
    code_lines = [line.rstrip() for line in code_text.splitlines() if line.strip()]
    if not code_lines:
        return []

    line_matched = []
    for shape in available:
        if id(shape) in assigned_ids or _is_heading_shape(shape):
            continue
        shape_lines = [line.rstrip() for line in _shape_text(shape).splitlines() if line.strip()]
        if any(
            _text_matches_code_line(shape_line, code_line)
            for shape_line in shape_lines
            for code_line in code_lines
        ):
            line_matched.append(shape)

    if not line_matched:
        min_score = max(40, sum(len(line) for line in code_lines) // 4)
        line_matched = [
            shape
            for shape in available
            if id(shape) not in assigned_ids
            and _code_block_match_score(_shape_text(shape), code_text) >= min_score
        ]

    if not line_matched:
        return []

    bounds = _bounds_for_shapes(line_matched)
    expanded = list(line_matched)
    for shape in available:
        if id(shape) in assigned_ids or shape in expanded:
            continue
        if _shape_in_bounds(shape, bounds, padding=120000) and (
            _shape_uses_monospace(shape)
            or _shape_has_code_fill(shape)
            or any(
                _text_matches_code_line(line.rstrip(), code_line)
                for line in _shape_text(shape).splitlines()
                for code_line in code_lines
            )
        ):
            expanded.append(shape)

    matched_lines = set()
    for shape in expanded:
        for shape_line in _shape_text(shape).splitlines():
            for code_line in code_lines:
                if _text_matches_code_line(shape_line.rstrip(), code_line):
                    matched_lines.add(code_line)

    required = max(1, min(len(code_lines), max(2, len(code_lines) // 3)))
    if len(matched_lines) < required:
        return []

    return expanded


def _code_block_step_numbers(slide_md: str) -> list[int | None]:
    step: int | None = None
    block_steps: list[int | None] = []
    in_fence = False

    for line in slide_md.splitlines():
        step_match = re.match(r"\*\*Step (\d+):\*\*", line.strip())
        if step_match:
            step = int(step_match.group(1))
        stripped = line.strip()
        if stripped.startswith("```"):
            if not in_fence:
                block_steps.append(step)
                in_fence = True
            else:
                in_fence = False

    return block_steps


def _find_step_label_shape(slide, step_number: int):
    pattern = re.compile(rf"^Step {step_number}:", re.IGNORECASE)
    for shape in slide.shapes:
        if not shape.has_text_frame or _is_chrome_shape(shape) or _is_heading_shape(shape):
            continue
        if pattern.match(_shape_text(shape).strip()):
            return shape
    return None


def _insert_top_for_code_block(
    slide,
    slide_md: str | None,
    block_index: int,
    *,
    slide_height: int,
    default_top: int,
) -> int:
    if slide_md is None:
        return default_top

    step_numbers = _code_block_step_numbers(slide_md)
    if block_index >= len(step_numbers):
        return default_top

    step_number = step_numbers[block_index]
    if step_number is None:
        return default_top

    step_shape = _find_step_label_shape(slide, step_number)
    if step_shape is None:
        return default_top

    _fit_textbox_height(step_shape)
    return int(step_shape.top) + int(step_shape.height) + int(CONTENT_GAP)


def _exclusive_code_shape_assignments(
    available,
    code_blocks: list[str],
) -> list[list]:
    assignments: list[list] = [[] for _ in code_blocks]
    scored: list[tuple] = []

    for shape in available:
        if _is_heading_shape(shape):
            continue
        shape_text = _shape_text(shape)
        if not shape_text:
            continue
        best_index = -1
        best_score = 0
        for index, code_text in enumerate(code_blocks):
            score = _code_block_match_score(shape_text, code_text)
            if score > best_score:
                best_score = score
                best_index = index
        if best_index >= 0 and best_score >= 8:
            scored.append((best_score, best_index, shape))

    scored.sort(key=lambda item: item[0], reverse=True)
    used_ids: set[int] = set()
    for _, block_index, shape in scored:
        if id(shape) in used_ids:
            continue
        assignments[block_index].append(shape)
        used_ids.add(id(shape))

    for index, code_text in enumerate(code_blocks):
        if assignments[index]:
            continue
        fallback = _find_fallback_code_shapes(available, used_ids, code_text)
        for shape in fallback:
            if id(shape) not in used_ids:
                assignments[index].append(shape)
                used_ids.add(id(shape))

    return assignments


def _code_block_match_threshold(code_text: str) -> int:
    code_lines = [line.rstrip() for line in code_text.splitlines() if line.strip()]
    if not code_lines:
        return 40
    total = sum(max(len(line), 8) for line in code_lines)
    return max(40, total // 2)


def _code_shape_matches_any_block(shape_text: str, code_blocks: list[str]) -> bool:
    if not shape_text.strip():
        return False
    for block in code_blocks:
        if _code_block_match_score(shape_text, block) >= _code_block_match_threshold(block):
            return True
        if _normalize_text(shape_text) == _normalize_text(block):
            return True
    return False


def _convert_shape_to_body_text(shape, *, min_pt: float = MIN_BODY_FONT_PT) -> None:
    if not shape.has_text_frame:
        return
    try:
        shape.fill.background()
    except Exception:
        pass
    try:
        shape.line.fill.background()
    except Exception:
        pass
    for paragraph in shape.text_frame.paragraphs:
        _apply_body_paragraph_spacing(paragraph)
        for run in paragraph.runs:
            run.font.name = BODY_FONT
            run.font.size = Pt(min_pt)
            run.font.color.rgb = BODY_BLACK


def _body_text_on_slide(slide) -> str:
    parts: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame or _is_chrome_shape(shape) or _is_heading_shape(shape):
            continue
        if _is_code_shape(shape):
            continue
        text = _shape_text(shape)
        if text:
            parts.append(text)
    return _normalize_text("\n".join(parts))


def _demote_false_code_shapes(slide, code_blocks: list[str]) -> None:
    body_text = _body_text_on_slide(slide)
    for shape in list(slide.shapes):
        if not _is_code_shape(shape) or _is_inline_monospace_shape(shape):
            continue
        shape_text = _shape_text(shape)
        if _code_shape_matches_any_block(shape_text, code_blocks):
            continue
        norm_shape = _normalize_text(shape_text)
        if norm_shape and norm_shape in body_text:
            sp = shape._element
            sp.getparent().remove(sp)
            continue
        _convert_shape_to_body_text(shape)


def _enforce_code_block_count(slide, code_blocks: list[str]) -> None:
    code_shapes = [
        shape
        for shape in slide.shapes
        if _is_code_shape(shape) and not _is_inline_monospace_shape(shape)
    ]
    if len(code_shapes) <= len(code_blocks):
        return

    non_matching = [
        shape
        for shape in code_shapes
        if not _code_shape_matches_any_block(_shape_text(shape), code_blocks)
    ]
    if non_matching:
        for shape in non_matching:
            _convert_shape_to_body_text(shape)
        code_shapes = [
            shape
            for shape in slide.shapes
            if _is_code_shape(shape) and not _is_inline_monospace_shape(shape)
        ]
        if len(code_shapes) <= len(code_blocks):
            return

    scored: list[tuple[int, object]] = []
    for shape in code_shapes:
        best_score = max(
            _code_block_match_score(_shape_text(shape), block) for block in code_blocks
        )
        scored.append((best_score, shape))
    scored.sort(key=lambda item: item[0])
    excess = len(code_shapes) - len(code_blocks)
    _remove_shapes([shape for _, shape in scored[:excess]])


def _dedupe_extra_code_shapes(slide, code_blocks: list[str]) -> None:
    code_shapes = [
        shape
        for shape in slide.shapes
        if _is_code_shape(shape) and not _is_inline_monospace_shape(shape)
    ]
    by_block: dict[int, list[tuple[int, object]]] = {}
    orphans: list = []

    for shape in code_shapes:
        text = _shape_text(shape)
        best_index = -1
        best_score = 0
        for index, block in enumerate(code_blocks):
            score = _code_block_match_score(text, block)
            if score > best_score:
                best_score = score
                best_index = index
        if best_index >= 0 and best_score >= _code_block_match_threshold(code_blocks[best_index]):
            by_block.setdefault(best_index, []).append((best_score, shape))
        else:
            orphans.append(shape)

    to_remove = list(orphans)
    for scored_shapes in by_block.values():
        scored_shapes.sort(key=lambda item: item[0], reverse=True)
        to_remove.extend(shape for _, shape in scored_shapes[1:])

    if to_remove:
        _remove_shapes(to_remove)


def _strip_code_fragments_from_body_shapes(slide, code_blocks: list[str]) -> None:
    if not code_blocks:
        return
    code_norm = _normalize_text("\n".join(code_blocks))
    for shape in slide.shapes:
        if not shape.has_text_frame or _is_chrome_shape(shape) or _is_heading_shape(shape):
            continue
        if _is_code_shape(shape):
            continue
        text = _shape_text(shape)
        if not text:
            continue
        if re.search(r"^Step \d+:", text.strip(), re.MULTILINE | re.IGNORECASE):
            continue
        norm_text = _normalize_text(text)
        if len(norm_text) < 8:
            continue
        if norm_text in code_norm and len(norm_text) < len(code_norm) * 0.35:
            sp = shape._element
            sp.getparent().remove(sp)


def _code_lines(code_text: str) -> list[str]:
    if not code_text:
        return [""]
    return code_text.split("\n")


def _format_code_text_frame(text_frame, code_text: str) -> None:
    text_frame.clear()
    text_frame.word_wrap = False
    text_frame.auto_size = MSO_AUTO_SIZE.NONE
    text_frame.vertical_anchor = MSO_ANCHOR.TOP

    for index, line in enumerate(_code_lines(code_text)):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(0)
        paragraph.line_spacing = Pt(CODE_LINE_HEIGHT_PT)
        run = paragraph.add_run()
        run.text = line
        run.font.name = CODE_FONT
        run.font.size = Pt(CODE_FONT_PT)
        run.font.color.rgb = CODE_TEXT


def _code_block_size(code_text: str, max_width: int | None = None) -> tuple[int, int]:
    lines = _code_lines(code_text)
    line_count = len(lines)
    max_chars = max(len(line) for line in lines)

    line_height = int(Pt(CODE_LINE_HEIGHT_PT).emu)
    char_width = int(Pt(CODE_CHAR_WIDTH_PT).emu)
    frame_pad = int(CODE_INSET) * 2
    buffer = int(Pt(CODE_SIZE_BUFFER_PT).emu)

    width = max_chars * char_width + frame_pad + buffer
    height = line_count * line_height + frame_pad + buffer

    if max_width is not None:
        width = min(width, max_width)

    return width, height


def _insert_code_block(
    slide,
    code_text: str,
    bounds: tuple[int, int, int, int],
    *,
    slide_width: int | None = None,
) -> None:
    left, top, _, _ = bounds
    max_width = int(slide_width * 0.9) if slide_width else None
    width, height = _code_block_size(code_text, max_width=max_width)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.fill.solid()
    textbox.fill.fore_color.rgb = CODE_BG
    textbox.line.color.rgb = CODE_BORDER
    textbox.line.width = Pt(0.75)

    frame = textbox.text_frame
    frame.margin_left = CODE_INSET
    frame.margin_right = CODE_INSET
    frame.margin_top = CODE_INSET
    frame.margin_bottom = CODE_INSET
    _format_code_text_frame(frame, code_text)
    return textbox


def _rebuild_slide_code_blocks(
    slide,
    code_blocks: list[str],
    *,
    slide_width: int | None = None,
    slide_height: int | None = None,
    slide_md: str | None = None,
) -> int:
    if not code_blocks:
        return 0

    assert slide_width is not None and slide_height is not None
    available = _content_text_shapes(slide)
    assignments = _exclusive_code_shape_assignments(available, code_blocks)

    for matched in assignments:
        if matched:
            _remove_shapes(matched)

    default_top = _estimate_code_insert_top(slide, slide_height)
    left = _content_left_margin(slide_width)
    insert_cursor = default_top
    rebuilt = 0

    for index, code_text in enumerate(code_blocks):
        top = _insert_top_for_code_block(
            slide,
            slide_md,
            index,
            slide_height=slide_height,
            default_top=insert_cursor,
        )
        bounds = (left, top, left, top)
        inserted = _insert_code_block(slide, code_text, bounds, slide_width=slide_width)
        insert_cursor = int(inserted.top) + int(inserted.height) + int(CONTENT_GAP)
        rebuilt += 1

    _dedupe_extra_code_shapes(slide, code_blocks)
    _demote_false_code_shapes(slide, code_blocks)
    _enforce_code_block_count(slide, code_blocks)
    _strip_code_fragments_from_body_shapes(slide, code_blocks)
    return rebuilt


def rebuild_code_blocks_from_markdown(prs: Presentation, md_path: Path) -> int:
    markdown = md_path.read_text(encoding="utf-8")
    md_slides = split_marp_slides(markdown)
    rebuilt = 0
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    search_from = 0

    for slide_md in md_slides:
        code_blocks = extract_fenced_code_blocks(slide_md)
        if not code_blocks:
            continue
        slide_index = _find_pptx_slide_index(prs, slide_md, start_index=search_from)
        if slide_index is None:
            continue
        search_from = slide_index + 1
        rebuilt += _rebuild_slide_code_blocks(
            prs.slides[slide_index],
            code_blocks,
            slide_width=slide_width,
            slide_height=slide_height,
            slide_md=slide_md,
        )

    return rebuilt


def _content_layout_shapes(slide) -> list:
    shapes = []
    for shape in slide.shapes:
        if _is_chrome_shape(shape):
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            shapes.append(shape)
        elif shape.has_text_frame and _shape_text(shape):
            shapes.append(shape)
    return shapes


def _horizontal_overlap(first, second, gap: int = 0) -> bool:
    left_a, _, right_a, _ = _shape_bounds(first)
    left_b, _, right_b, _ = _shape_bounds(second)
    return not (right_a + gap <= left_b or right_b + gap <= left_a)


def _shape_max_font_pt(shape, default: float = MIN_BODY_FONT_PT) -> float:
    max_pt = default
    if not shape.has_text_frame:
        return default
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                max_pt = max(max_pt, run.font.size.pt)
    if _is_heading_shape(shape):
        max_pt = max(max_pt, HEADING_FONT_PT)
    return max_pt


def _estimate_wrapped_line_count(shape) -> int:
    if not shape.has_text_frame:
        return 0
    font_pt = _shape_max_font_pt(shape)
    char_width = int(Pt(max(font_pt * 0.52, 7)).emu)
    frame_inset = int(Emu(80000))
    usable_width = max(int(shape.width) - frame_inset * 2, char_width)
    chars_per_line = max(8, usable_width // char_width)
    lines = 0
    for paragraph in shape.text_frame.paragraphs:
        text = (paragraph.text or "".join(run.text for run in paragraph.runs)).replace("\r", "")
        if not text.strip():
            continue
        for raw in text.split("\n"):
            chunk = raw.strip()
            if not chunk:
                continue
            lines += max(1, (len(chunk) + chars_per_line - 1) // chars_per_line)
    return lines


def _fit_textbox_height(shape) -> None:
    if not shape.has_text_frame or _is_chrome_shape(shape):
        return
    if _is_code_shape(shape):
        line_count = max(1, len(_code_lines(_shape_text(shape))))
        line_height = int(Pt(CODE_LINE_HEIGHT_PT).emu)
        frame_pad = int(CODE_INSET) * 2
        buffer = int(Pt(CODE_SIZE_BUFFER_PT).emu)
        shape.height = line_count * line_height + frame_pad + buffer
        return

    line_count = _estimate_wrapped_line_count(shape)
    if line_count == 0:
        return
    font_pt = _shape_max_font_pt(shape)
    line_height = int(Pt(max(font_pt * BODY_LINE_SPACING, MIN_BODY_FONT_PT)).emu)
    padding = int(Emu(60000 if _is_heading_shape(shape) else 40000))
    shape.height = line_count * line_height + padding


def _trim_oversized_body_textbox(shape) -> None:
    _fit_textbox_height(shape)


def _trim_oversized_body_textboxes(slide) -> None:
    for shape in slide.shapes:
        _fit_textbox_height(shape)


def _layout_shapes(slide) -> list:
    return [
        shape
        for shape in _content_layout_shapes(slide)
        if not _is_code_shape(shape)
    ]


def _rects_overlap(
    left_a: int,
    top_a: int,
    right_a: int,
    bottom_a: int,
    left_b: int,
    top_b: int,
    right_b: int,
    bottom_b: int,
    gap: int = 0,
) -> bool:
    return not (
        right_a + gap <= left_b
        or right_b + gap <= left_a
        or bottom_a + gap <= top_b
        or bottom_b + gap <= top_a
    )


def _shape_pair_overlaps(first, second, gap: int = 0) -> bool:
    a = _shape_bounds(first)
    b = _shape_bounds(second)
    return _rects_overlap(a[0], a[1], a[2], a[3], b[0], b[1], b[2], b[3], gap)


def _slide_has_layout_issues(slide, slide_height: int) -> bool:
    shapes = _content_layout_shapes(slide)
    footer_line = int(slide_height * FOOTER_SAFE_RATIO)
    if any(_shape_bounds(shape)[3] > footer_line for shape in shapes):
        return True
    for index, first in enumerate(shapes):
        for second in shapes[index + 1 :]:
            if not _horizontal_overlap(first, second):
                continue
            gap = CONTENT_GAP
            if _is_heading_shape(first) and _is_heading_shape(second):
                gap = Emu(20000)
            if _shape_pair_overlaps(first, second, gap):
                return True
    return False


def _markdown_content_sequence(slide_md: str) -> list[tuple[str, str]]:
    sequence: list[tuple[str, str]] = []
    in_fence = False
    fence_lines: list[str] = []

    for line in slide_md.splitlines():
        stripped = line.strip()
        if stripped == "---":
            break
        if (
            stripped.startswith("## ")
            or stripped.startswith("# ")
            or stripped.startswith("<!--")
        ):
            continue
        if stripped.startswith("```"):
            if not in_fence:
                in_fence = True
                fence_lines = []
            else:
                in_fence = False
                sequence.append(("code", "\n".join(fence_lines)))
            continue
        if in_fence:
            fence_lines.append(line)
            continue
        if not stripped:
            continue

        step_parts = list(re.finditer(r"\*\*Step (\d+):\*\*\s*([^*]*)", stripped))
        if step_parts:
            for match in step_parts:
                label = f"Step {match.group(1)}: {match.group(2).strip().rstrip(' ·')}"
                sequence.append(("step", label))
            continue

        if stripped.startswith("- "):
            sequence.append(("body", stripped[2:].strip()))
        elif stripped.startswith("**") and stripped.endswith("**"):
            sequence.append(("body", stripped.strip("*").strip()))
        elif not stripped.startswith("|"):
            sequence.append(("body", stripped))

    return sequence


def _sequence_item_match_score(shape, kind: str, text: str) -> int:
    shape_text = _shape_text(shape)
    if not shape_text:
        return 0
    if kind == "code":
        if not _is_code_shape(shape):
            return 0
        return _code_block_match_score(shape_text, text)

    if _is_code_shape(shape):
        return 0

    norm_item = _normalize_text(text)
    norm_shape = _normalize_text(shape_text)
    if not norm_item:
        return 0
    if norm_item in norm_shape:
        return len(norm_item) + 100
    if norm_shape in norm_item:
        return len(norm_shape) + 80
    step_match = re.match(r"Step (\d+):", text, re.I)
    if step_match and re.search(rf"Step {step_match.group(1)}:", shape_text, re.I):
        return 120
    prefix = norm_item[: min(24, len(norm_item))]
    if prefix and shape_text.strip().lower().startswith(text[: min(24, len(text))].lower()):
        return 60
    return 0


def _ordered_content_shapes(slide, slide_md: str | None) -> list:
    content = [
        shape
        for shape in _content_layout_shapes(slide)
        if not _is_heading_shape(shape)
    ]
    if not slide_md:
        content.sort(key=lambda shape: (int(shape.top), int(shape.left)))
        return content

    sequence = _markdown_content_sequence(slide_md)
    if not sequence:
        content.sort(key=lambda shape: (int(shape.top), int(shape.left)))
        return content

    ordered: list = []
    used_ids: set[int] = set()
    for kind, text in sequence:
        best_shape = None
        best_score = 0
        for shape in content:
            shape_id = id(shape._element)
            if shape_id in used_ids:
                continue
            score = _sequence_item_match_score(shape, kind, text)
            if score > best_score:
                best_score = score
                best_shape = shape
        if best_shape is not None and best_score > 0:
            ordered.append(best_shape)
            used_ids.add(id(best_shape._element))

    for shape in sorted(content, key=lambda item: (int(item.top), int(item.left))):
        if id(shape._element) not in used_ids:
            ordered.append(shape)
    return ordered


def _resolve_content_flow(
    slide,
    slide_height: int,
    *,
    content_order: list | None = None,
) -> None:
    footer_line = int(slide_height * FOOTER_SAFE_RATIO)
    heading_top = int(HEADING_TOP_CONTENT)

    if _is_course_agenda_slide(slide):
        return

    if _is_lesson_lead_slide(slide):
        return

    if _is_module_lead_slide(slide):
        return

    if _is_up_next_slide(slide):
        return

    all_shapes = _content_layout_shapes(slide)
    headings = [shape for shape in all_shapes if _is_heading_shape(shape)]
    if content_order is None:
        content = [shape for shape in all_shapes if not _is_heading_shape(shape)]
        content.sort(key=lambda shape: (int(shape.top), int(shape.left)))
    else:
        content = list(content_order)

    if headings:
        primary = min(headings, key=lambda shape: int(shape.top))
        primary.top = heading_top
        _fit_textbox_height(primary)
        cursor = int(primary.top) + int(primary.height) + HEADING_BODY_GAP
        stacked = sorted(headings, key=lambda shape: int(shape.top))
        for index in range(1, len(stacked)):
            previous = stacked[index - 1]
            current = stacked[index]
            _fit_textbox_height(current)
            min_top = _shape_bounds(previous)[3] + Emu(20000)
            if int(current.top) < min_top:
                current.top = min_top
            cursor = max(cursor, int(current.top) + int(current.height) + HEADING_BODY_GAP)
    else:
        cursor = heading_top

    for shape in content:
        _fit_textbox_height(shape)
        shape.top = cursor
        cursor = int(shape.top) + int(shape.height) + CONTENT_GAP

    for _ in range(12):
        shapes = _content_layout_shapes(slide)
        if not shapes:
            return
        max_bottom = max(_shape_bounds(shape)[3] for shape in shapes)
        if max_bottom <= footer_line:
            return
        shift = max_bottom - footer_line + Emu(40000)
        for shape in shapes:
            shape.top = max(0, int(shape.top) - shift)
        if headings:
            primary = min(headings, key=lambda shape: int(shape.top))
            cursor = int(primary.top) + int(primary.height) + HEADING_BODY_GAP
            for shape in content:
                _fit_textbox_height(shape)
                shape.top = cursor
                cursor = int(shape.top) + int(shape.height) + CONTENT_GAP


def _fit_oversized_code_blocks(slide, slide_height: int) -> None:
    footer_line = int(slide_height * FOOTER_SAFE_RATIO)

    for shape in slide.shapes:
        if not _is_code_shape(shape):
            continue

        max_bottom = footer_line - int(CONTENT_GAP)
        available = max_bottom - int(shape.top)
        if available <= 0 or int(shape.height) <= available:
            continue

        scale = available / int(shape.height)
        font_pt = max(11, CODE_FONT_PT * scale)
        line_pt = max(12, CODE_LINE_HEIGHT_PT * scale)
        line_count = max(1, len(_code_lines(_shape_text(shape))))

        for paragraph in shape.text_frame.paragraphs:
            paragraph.line_spacing = Pt(line_pt)
            for run in paragraph.runs:
                run.font.size = Pt(font_pt)

        frame_pad = int(CODE_INSET) * 2
        buffer = int(Pt(CODE_SIZE_BUFFER_PT).emu)
        line_height = int(Pt(line_pt).emu)
        shape.height = min(available, line_count * line_height + frame_pad + buffer)


def _shrink_table_fonts(shape, delta_pt: float) -> None:
    if shape.shape_type != MSO_SHAPE_TYPE.TABLE:
        return
    table = shape.table
    for row_idx in range(len(table.rows)):
        for col_idx in range(len(table.columns)):
            cell = table.cell(row_idx, col_idx)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    if run.font.size:
                        run.font.size = Pt(max(10, run.font.size.pt + delta_pt))


def _shrink_shape_fonts(shape, delta_pt: float) -> None:
    if not shape.has_text_frame or _is_code_shape(shape) or _is_heading_shape(shape):
        return
    shape.text_frame.word_wrap = True
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size and run.font.size.pt > MIN_LAYOUT_FONT_PT:
                run.font.size = Pt(max(MIN_LAYOUT_FONT_PT, run.font.size.pt + delta_pt))


def _nudge_shape_up(shape, amount: int) -> None:
    shape.top = max(0, int(shape.top) - amount)


def _fix_stacked_tables(slide) -> None:
    tables = [
        shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.TABLE
    ]
    if len(tables) < 2:
        return

    tables.sort(key=lambda shape: int(shape.top))
    for index in range(1, len(tables)):
        previous = tables[index - 1]
        current = tables[index]
        prev_bottom = int(previous.top) + int(previous.height)
        gap = Emu(120000)
        if int(current.top) < prev_bottom + gap:
            current.top = prev_bottom + gap


def _slide_has_non_table_body(slide) -> bool:
    for shape in slide.shapes:
        if _is_chrome_shape(shape) or _is_heading_shape(shape):
            continue
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            continue
        if _is_code_shape(shape):
            return True
        if shape.has_text_frame and _shape_text(shape).strip():
            return True
    return False


def _format_table_only_slide(slide, slide_height: int, slide_width: int) -> None:
    """Center a lone table when the slide has no body text or code blocks."""
    if _slide_has_non_table_body(slide):
        return

    tables = [
        shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.TABLE
    ]
    if len(tables) != 1:
        return

    heading_bottom = int(HEADING_TOP_CONTENT)
    headings = [shape for shape in slide.shapes if _is_heading_shape(shape)]
    if headings:
        primary = min(headings, key=lambda shape: int(shape.top))
        heading_bottom = int(primary.top) + int(primary.height) + int(HEADING_BODY_GAP)

    footer_line = int(slide_height * FOOTER_SAFE_RATIO)
    table = tables[0]
    table.left = max(0, int((slide_width - int(table.width)) // 2))
    available = footer_line - heading_bottom
    if available > int(table.height):
        table.top = heading_bottom + (available - int(table.height)) // 2


def _fix_slide_layout(
    slide,
    slide_height: int,
    slide_width: int,
    *,
    slide_md: str | None = None,
) -> None:
    _fix_stacked_tables(slide)
    _trim_oversized_body_textboxes(slide)

    initial_content = _ordered_content_shapes(slide, slide_md)

    _resolve_content_flow(slide, slide_height, content_order=initial_content)
    _trim_oversized_body_textboxes(slide)
    initial_content = _ordered_content_shapes(slide, slide_md)
    _resolve_content_flow(slide, slide_height, content_order=initial_content)
    _align_content_slide(slide, slide_width)
    _format_table_only_slide(slide, slide_height, slide_width)
    _fit_oversized_code_blocks(slide, slide_height)

    footer_line = int(slide_height * FOOTER_SAFE_RATIO)
    for _ in range(20):
        if not _slide_has_layout_issues(slide, slide_height):
            return

        _resolve_content_flow(slide, slide_height, content_order=initial_content)
        _align_content_slide(slide, slide_width)
        _format_table_only_slide(slide, slide_height, slide_width)
        _trim_oversized_body_textboxes(slide)
        _fit_oversized_code_blocks(slide, slide_height)
        initial_content = _ordered_content_shapes(slide, slide_md)

        for shape in _layout_shapes(slide):
            if _shape_bounds(shape)[3] > footer_line:
                _nudge_shape_up(shape, Emu(80000))
            if shape.has_text_frame:
                shape.text_frame.word_wrap = True
            _shrink_shape_fonts(shape, -1)
            _shrink_table_fonts(shape, -1)
            _trim_oversized_body_textbox(shape)


def fix_presentation_layout(prs: Presentation, source_md: Path | None = None) -> None:
    slide_height = int(prs.slide_height)
    slide_width = int(prs.slide_width)
    md_slides: list[str] = []
    if source_md is not None:
        md_slides = split_marp_slides(source_md.read_text(encoding="utf-8"))
    for idx, slide in enumerate(prs.slides):
        slide_md = md_slides[idx] if idx < len(md_slides) else None
        if _is_course_title_slide(slide) or _is_course_agenda_slide(slide):
            continue
        if _is_lesson_lead_slide(slide):
            _format_lesson_lead_slide(slide, slide_height, slide_width)
            continue
        if _is_module_lead_slide(slide):
            _format_module_lead_slide(slide, slide_height, slide_width)
            continue
        if _is_up_next_slide(slide):
            _format_up_next_slide(slide, slide_height, slide_width)
            continue
        _fix_slide_layout(slide, slide_height, slide_width, slide_md=slide_md)
    _format_course_intro_slides(prs)


def rebuild_tables_from_markdown(prs: Presentation, md_path: Path) -> int:
    markdown = md_path.read_text(encoding="utf-8")
    md_slides = split_marp_slides(markdown)
    rebuilt = 0
    slide_width = int(prs.slide_width)
    search_from = 0

    for slide_md in md_slides:
        tables = extract_markdown_tables(slide_md)
        if not tables:
            continue
        slide_index = _find_pptx_slide_index(prs, slide_md, start_index=search_from)
        if slide_index is None:
            continue
        search_from = slide_index + 1
        rebuilt += _rebuild_slide_tables(
            prs.slides[slide_index], tables, slide_width=slide_width
        )

    return rebuilt


def clean_pptx(
    input_path: Path,
    output_path: Path,
    *,
    mode: str = "minimal",
    dry_run: bool = False,
    source_md: Path | None = None,
) -> tuple[int, int, int, int, int, int]:
    prs = Presentation(str(input_path))
    slide_width = int(prs.slide_width)
    slide_height = int(prs.slide_height)
    removed = 0
    slides_touched = 0
    tables_rebuilt = 0
    code_blocks_rebuilt = 0
    chrome_slides_updated = 0
    body_groups_merged = 0

    for slide in prs.slides:
        to_remove = []

        for shape in slide.shapes:
            text = _shape_text(shape)
            rgb = _shape_fill_rgb(shape)
            full = _is_full_slide(shape, slide_width, slide_height)

            if mode == "minimal":
                if full and rgb == "FFFFFF":
                    to_remove.append(shape)
                continue

            # text-only: keep shapes with readable text; drop decoration meshes.
            if text:
                continue
            to_remove.append(shape)

        if to_remove:
            slides_touched += 1
            removed += len(to_remove)
            if not dry_run:
                for shape in to_remove:
                    sp = shape._element
                    sp.getparent().remove(sp)

        if mode == "text-only" and not dry_run:
            _set_slide_background(slide)
            for shape in slide.shapes:
                _scale_body_fonts(shape)
                _apply_text_colors(shape)

    if mode == "text-only" and source_md and not dry_run:
        tables_rebuilt = rebuild_tables_from_markdown(prs, source_md)
        code_blocks_rebuilt = rebuild_code_blocks_from_markdown(prs, source_md)
        md_slides = split_marp_slides(source_md.read_text(encoding="utf-8"))
        for slide in prs.slides:
            _absorb_inline_monospace_shapes(slide)
            _coalesce_same_line_body_fragments(slide)
        body_groups_merged = merge_body_textboxes(prs, md_slides=md_slides)
        for slide in prs.slides:
            _absorb_inline_monospace_shapes(slide)
            _coalesce_same_line_body_fragments(slide)
        body_groups_merged += merge_body_textboxes(prs, md_slides=md_slides)
        fix_presentation_layout(prs, source_md=source_md)
        _normalize_presentation_formatting(prs)
        fix_presentation_layout(prs, source_md=source_md)
        _format_course_intro_slides(prs)
        for slide in prs.slides:
            if _is_lesson_lead_slide(slide):
                _format_lesson_lead_slide(slide, slide_height, slide_width)
            elif _is_module_lead_slide(slide):
                _format_module_lead_slide(slide, slide_height, slide_width)
            elif _is_up_next_slide(slide):
                _format_up_next_slide(slide, slide_height, slide_width)
        _normalize_presentation_formatting(prs)
        chrome_slides_updated = apply_slide_chrome_from_markdown(prs, source_md)
        for slide in prs.slides:
            _fit_oversized_code_blocks(slide, slide_height)
        for slide in prs.slides:
            if _is_up_next_slide(slide):
                _format_up_next_slide(slide, slide_height, slide_width)

    if not dry_run:
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))

    return removed, slides_touched, tables_rebuilt, code_blocks_rebuilt, chrome_slides_updated, body_groups_merged


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("input", type=Path)
    parser.add_argument("-o", "--output", type=Path)
    parser.add_argument(
        "--mode",
        choices=("minimal", "text-only"),
        default="text-only",
        help="minimal removes white overlays; text-only keeps text boxes + flat background",
    )
    parser.add_argument("--dry-run", action="store_true")
    parser.add_argument(
        "--source-md",
        type=Path,
        help="Marp markdown source used to rebuild pipe tables as native PowerPoint tables",
    )
    args = parser.parse_args()

    if not args.input.exists():
        print(f"Not found: {args.input}", file=sys.stderr)
        return 1

    output = args.output or args.input.with_name(
        args.input.stem.replace(".cleaned", "") + ".cleaned.pptx"
    )

    removed, slides, tables_rebuilt, code_blocks_rebuilt, chrome_slides_updated, body_groups_merged = clean_pptx(
        args.input, output, mode=args.mode, dry_run=args.dry_run, source_md=args.source_md
    )
    action = "Would remove" if args.dry_run else "Removed"
    print(f"{action} {removed} shape(s) across {slides} slide(s) [{args.mode}]")
    if tables_rebuilt:
        print(f"Rebuilt {tables_rebuilt} native table(s) from markdown")
    if code_blocks_rebuilt:
        print(f"Rebuilt {code_blocks_rebuilt} styled code block(s) from markdown")
    if body_groups_merged:
        print(f"Merged fragmented body text into {body_groups_merged} text box(es)")
    if chrome_slides_updated:
        print(f"Updated header/footer chrome on {chrome_slides_updated} slide(s) from markdown")
    if not args.dry_run:
        print(f"Saved: {output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
